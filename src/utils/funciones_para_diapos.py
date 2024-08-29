from pptx.util import Pt
from pptx.dml.color import RGBColor
from src.config.config import DEFAULT_FONT_NAME, DEFAULT_CHAR_WIDTH
from PIL import Image as PilImage
from pptx.util import Inches

def add_resized_image_to_slide(slide, img_path, maxima_anchura_percent, prs):
    """
    Abre una imagen, la convierte a RGBA, la redimensiona según las dimensiones máximas de la diapositiva y la añade a la diapositiva.

    :param slide: Objeto Slide de python-pptx donde se añadirá la imagen.
    :param img_path: Ruta de la imagen a añadir.
    :param positions_initial: Diccionario que contiene las posiciones iniciales y dimensiones máximas para la imagen.
    :param prs: Objeto Presentation de python-pptx que contiene las dimensiones de la diapositiva.
    """
    # Abrir y convertir la imagen
    imagen_orchestra = PilImage.open(img_path).convert("RGBA")
    orchestra_image_width, orchestra_image_height = imagen_orchestra.size

    # Convertir dimensiones de píxeles a pulgadas
    puntos_por_pulgada = 80
    orchestra_image_width = orchestra_image_width / puntos_por_pulgada
    orchestra_image_height = orchestra_image_height / puntos_por_pulgada

    # Convertir dimensiones de pulgadas a EMUs
    emus_por_pulgada = 914400
    orchestra_image_width = orchestra_image_width * emus_por_pulgada
    orchestra_image_height = orchestra_image_height * emus_por_pulgada

    # Obtener las dimensiones máximas permitidas para la imagen
    maximum_in_width = maxima_anchura_percent * prs.slide_width
    maximum_in_height = prs.slide_height

    # Calcular el factor de redimensionamiento basado en la altura
    resize_factor = maximum_in_height / orchestra_image_height

    # Ajustar el factor de redimensionamiento si el ancho redimensionado excede el máximo permitido
    if (resize_factor * orchestra_image_width) > maximum_in_width:
        resize_factor = maximum_in_width / orchestra_image_width

    # Calcular las nuevas dimensiones redimensionadas
    new_width = int(orchestra_image_width * resize_factor)
    new_height = int(orchestra_image_height * resize_factor)

    # Añadir la imagen redimensionada a la diapositiva, alineada al fondo
    image = slide.shapes.add_picture(img_path, 0, prs.slide_height - new_height, width=new_width, height=new_height)

    return image


def adjust_text_size(text_frame, max_width_cm, max_font_size=100, min_font_size=10, fuente=DEFAULT_FONT_NAME):

    # Convertir las medidas de cm a puntos (pt)
    max_width_pt = max_width_cm / 11500

    font_size = max_font_size
    text = ''.join(run.text for paragraph in text_frame.paragraphs for run in paragraph.runs)
    text_length = len(text)

    while font_size >= min_font_size:
        # Aplicar el tamaño de la fuente
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                run.font.name = fuente

        # Cálculo del ancho estimado del texto basado en el tamaño de la fuente y la longitud del texto
        char_width = (DEFAULT_CHAR_WIDTH / 100) * font_size
        estimated_text_width = char_width * text_length

        # Si el ancho estimado del texto es menor o igual al ancho máximo permitido, se rompe el bucle
        if estimated_text_width <= max_width_pt:
            break

        font_size -= 1

    # Aplicar el tamaño de fuente final
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.name = fuente


def add_text_to_slide(slide, full_text, posicion, offset, tamano_fuente, font_name, font_color_rgb, is_bold,
                      has_shadow, extra_paragraph_text=None, extra_run_text=None,
                      extra_paragraph_settings=None, extra_run_settings=None,
                      border_color_rgb=None, border_width_pt=None):
    """
    Añade texto con sombra en una diapositiva, si se especifica. También permite agregar un segundo párrafo y run opcionalmente, ambos con sombras si se especifica.
    Además, permite añadir un borde alrededor del cuadro de texto.

    :param slide: Objeto Slide de python-pptx.
        El objeto de la diapositiva donde se va a añadir el texto.

    :param full_text: str
        Texto completo que se añadirá en la diapositiva (puede ser una concatenación de varios elementos, como initial_text y orchestra_value).

    :param posicion: list or tuple
        Lista o tupla con las posiciones [left, top, width, height] donde se ubicará el texto en la diapositiva.

    :param offset: int or float
        Desplazamiento en píxeles para las sombras, determinando cuánto se desplaza la sombra respecto al texto original.

    :param tamano_fuente: int or float
        Tamaño de la fuente que se aplicará al texto.

    :param font_name: str
        Nombre de la fuente que se utilizará para el texto.

    :param font_color_rgb: tuple
        Tupla con los valores RGB (rojo, verde, azul) que definen el color de la fuente, por ejemplo, (255, 255, 255) para blanco.

    :param is_bold: bool
        Booleano que indica si la fuente debe estar en negrita (True) o no (False).

    :param has_shadow: bool
        Booleano que indica si el texto debe tener sombra (True) o no (False). Si es True, se agregan las sombras al texto.

    :param extra_paragraph_text: str, optional
        Texto adicional para un segundo párrafo opcional.

    :param extra_run_text: str, optional
        Texto adicional para un segundo run opcional.

    :param extra_paragraph_settings: dict, optional
        Configuración para el segundo párrafo opcional, incluyendo 'font_name', 'tamano_fuente', 'font_color_rgb', 'is_bold', y 'is_italic'.

    :param extra_run_settings: dict, optional
        Configuración para el segundo run opcional, incluyendo 'font_name', 'tamano_fuente', 'font_color_rgb', 'is_bold', y 'is_italic'.

    :param border_color_rgb: tuple, optional
        Tupla con los valores RGB para definir el color del borde (por ejemplo, (0, 0, 0) para negro).

    :param border_width_pt: int or float, optional
        Grosor del borde en puntos (por ejemplo, 2 para un borde de 2 pt de grosor).
    """
    # Añadir las sombras si has_shadow es True para el texto principal
    if has_shadow:
        for x_offset in [-offset, offset]:
            for y_offset in [-offset, offset]:
                shadow = slide.shapes.add_textbox(
                    posicion[0] + x_offset,
                    posicion[1] + y_offset,
                    posicion[2],
                    posicion[3]
                )
                shadow_frame = shadow.text_frame
                shadow_paragraph = shadow_frame.paragraphs[0]
                shadow_run = shadow_paragraph.add_run()

                shadow_run.text = full_text
                shadow_run.font.size = Pt(tamano_fuente)
                shadow_run.font.color.rgb = RGBColor(0, 0, 0)  # Color negro para la sombra
                shadow_run.font.bold = is_bold
                shadow_run.font.name = font_name
                if extra_run_text:
                    extra_run = shadow_paragraph.add_run()
                    extra_run.text = extra_run_text

                    if extra_run_settings:
                        extra_run.font.name = extra_run_settings.get('font_name', font_name)
                        extra_run.font.size = Pt(extra_run_settings.get('tamano_fuente', tamano_fuente))
                        extra_run.font.color.rgb = RGBColor(0, 0, 0)
                        extra_run.font.bold = extra_run_settings.get('is_bold', is_bold)
                        extra_run.font.italic = extra_run_settings.get('is_italic', False)
                if extra_paragraph_text:
                    extra_paragraph = shadow_frame.add_paragraph()
                    extra_paragraph.text = extra_paragraph_text

                    if extra_paragraph_settings:
                        extra_paragraph.font.name = extra_paragraph_settings.get('font_name', font_name)
                        extra_paragraph.font.size = Pt(extra_paragraph_settings.get('tamano_fuente', tamano_fuente))
                        extra_paragraph.font.color.rgb = RGBColor(0, 0, 0)
                        extra_paragraph.font.bold = extra_paragraph_settings.get('is_bold', is_bold)
                        extra_paragraph.font.italic = extra_paragraph_settings.get('is_italic', False)

    # Añadir el texto principal
    text = slide.shapes.add_textbox(
        posicion[0],
        posicion[1],
        posicion[2],
        posicion[3]
    )
    title_frame = text.text_frame
    title_paragraph1 = title_frame.paragraphs[0]
    run_orquesta = title_paragraph1.add_run()

    run_orquesta.text = full_text
    run_orquesta.font.size = Pt(tamano_fuente)
    run_orquesta.font.color.rgb = RGBColor(*font_color_rgb)  # Color definido por el usuario
    run_orquesta.font.bold = is_bold
    run_orquesta.font.name = font_name

    # # Añadir sombra al segundo párrafo si se proporciona
    # if extra_paragraph_text and has_shadow:
    #     for x_offset in [-offset, offset]:
    #         for y_offset in [-offset, offset]:
    #             shadow_paragraph = slide.shapes.add_textbox(
    #                 posicion[0] + x_offset,
    #                 posicion[1] + y_offset + 100000,  # Añadir un desplazamiento vertical
    #                 posicion[2],
    #                 posicion[3]
    #             )
    #             shadow_frame = shadow_paragraph.text_frame
    #             shadow_paragraph = shadow_frame.add_paragraph()
    #             shadow_paragraph.text = extra_paragraph_text
    #
    #             if extra_paragraph_settings:
    #                 shadow_paragraph.font.size = Pt(extra_paragraph_settings.get('tamano_fuente', tamano_fuente))
    #                 shadow_paragraph.font.color.rgb = RGBColor(0, 0, 0)
    #                 shadow_paragraph.font.bold = extra_paragraph_settings.get('is_bold', is_bold)
    #                 shadow_paragraph.font.name = extra_paragraph_settings.get('font_name', font_name)
    #                 shadow_paragraph.font.italic = extra_paragraph_settings.get('is_italic', False)

    # Añadir el segundo párrafo si se proporciona
    if extra_paragraph_text:
        extra_paragraph = title_frame.add_paragraph()
        extra_paragraph.text = extra_paragraph_text

        if extra_paragraph_settings:
            extra_paragraph.font.name = extra_paragraph_settings.get('font_name', font_name)
            extra_paragraph.font.size = Pt(extra_paragraph_settings.get('tamano_fuente', tamano_fuente))
            extra_paragraph.font.color.rgb = RGBColor(*extra_paragraph_settings.get('font_color_rgb', font_color_rgb))
            extra_paragraph.font.bold = extra_paragraph_settings.get('is_bold', is_bold)
            extra_paragraph.font.italic = extra_paragraph_settings.get('is_italic', False)

    # # Añadir sombra al segundo run si se proporciona
    # if extra_run_text and has_shadow:
    #     for x_offset in [-offset, offset]:
    #         for y_offset in [-offset, offset]:
    #             shadow_run = slide.shapes.add_textbox(
    #                 posicion[0] + x_offset,
    #                 posicion[1] + y_offset,
    #                 posicion[2],
    #                 posicion[3]
    #             )
    #             shadow_frame = shadow_run.text_frame
    #             shadow_paragraph = shadow_frame.paragraphs[0]
    #             shadow_run = shadow_paragraph.add_run()
    #
    #             shadow_run.text = extra_run_text
    #             if extra_run_settings:
    #                 shadow_run.font.size = Pt(extra_run_settings.get('tamano_fuente', tamano_fuente))
    #                 shadow_run.font.color.rgb = RGBColor(0, 0, 0)
    #                 shadow_run.font.bold = extra_run_settings.get('is_bold', is_bold)
    #                 shadow_run.font.name = extra_run_settings.get('font_name', font_name)
    #                 shadow_run.font.italic = extra_run_settings.get('is_italic', False)

    # Añadir el segundo run si se proporciona
    if extra_run_text:
        extra_run = title_paragraph1.add_run()
        extra_run.text = extra_run_text

        if extra_run_settings:
            extra_run.font.name = extra_run_settings.get('font_name', font_name)
            extra_run.font.size = Pt(extra_run_settings.get('tamano_fuente', tamano_fuente))
            extra_run.font.color.rgb = RGBColor(*extra_run_settings.get('font_color_rgb', font_color_rgb))
            extra_run.font.bold = extra_run_settings.get('is_bold', is_bold)
            extra_run.font.italic = extra_run_settings.get('is_italic', False)

    # Añadir borde al cuadro de texto si se especifica
    if border_color_rgb and border_width_pt:
        text.line.color.rgb = RGBColor(*border_color_rgb)  # Color del borde
        text.line.width = Pt(border_width_pt)  # Grosor del borde


def transformar_posiciones(diaspo, posiciones_dic):
    # Obtén las dimensiones de la diapositiva
    ancho_diapositiva = diaspo.slide_width
    alto_diapositiva = diaspo.slide_height

    pos_transformada = {}

    for clave, valores in posiciones_dic.items():
        if isinstance(valores, dict):
            # Calcular distancia al borde izquierdo
            distancia_borde_izquierdo = valores.get("left")
            if distancia_borde_izquierdo is None and "right" in valores and "width" in valores:
                distancia_borde_izquierdo = ancho_diapositiva - valores["right"] - valores["width"]

            # Calcular distancia al borde superior
            distancia_borde_superior = valores.get("top")
            if distancia_borde_superior is None and "bottom" in valores and "height" in valores:
                distancia_borde_superior = alto_diapositiva - valores["bottom"] - valores["height"]

            # Calcular ancho
            ancho = valores.get("width")
            if ancho is None and "left" in valores and "right" in valores:
                ancho = ancho_diapositiva - valores["left"] - valores["right"]

            # Calcular altura
            altura = valores.get("height")
            if altura is None and "top" in valores and "bottom" in valores:
                altura = alto_diapositiva - valores["top"] - valores["bottom"]

            # Solo almacena la transformación si todos los valores son válidos
            if None not in (distancia_borde_izquierdo, distancia_borde_superior, ancho, altura):
                pos_transformada[clave] = [
                    distancia_borde_izquierdo,
                    distancia_borde_superior,
                    ancho,
                    altura
                ]
            else:
                # Si no se cumplen los valores, deja la entrada tal cual
                pos_transformada[clave] = valores
        else:
            # Para valores que no son diccionarios o que no contienen las claves necesarias, se mantienen igual
            pos_transformada[clave] = valores

    return pos_transformada