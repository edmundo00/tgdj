import re
import os
import pandas as pd
import pygame
import ftfy
import num2words as nw
from unidecode import unidecode
from pptx.util import Pt
from src.config.config import *
from src.constants.enums import TagLabels
from mutagen.flac import FLAC
from mutagen.mp4 import MP4
from mutagen.mp3 import MP3
from mutagen.easyid3 import EasyID3
from datetime import datetime
import difflib
from tinytag import TinyTag

def convertir_segundos(segundos, formato='x\'x\'\''):
    """
    Convierte segundos con decimales a un formato de tiempo especificado.

    Args:
        segundos (float): Tiempo en segundos.
        formato (str): Formato de salida. Puede ser:
                       - "x'x''": Formato de minutos y segundos (ej. 2'6'').
                       - "x minutos y x segundos": Texto completo (ej. 2 minutos y 6 segundos).
                       - "x min x sec": Abreviado (ej. 2 min 6 sec).
                       - "h:m:s": Formato reloj (ej. 1:02:06).

    Returns:
        str: Tiempo formateado según el estilo especificado.
    """
    horas = int(segundos // 3600)  # División entera para obtener las horas
    minutos = int((segundos % 3600) // 60)  # Minutos restantes
    segundos_restantes = round(segundos % 60)  # Segundos restantes redondeados

    if formato == "x'x''":
        if horas > 0:
            return f"{horas}h {minutos}'{segundos_restantes}''"
        else:
            return f"{minutos}'{segundos_restantes}''"
    elif formato == "x minutos y x segundos":
        if horas > 0:
            return f"{horas} horas, {minutos} minutos y {segundos_restantes} segundos"
        else:
            return f"{minutos} minutos y {segundos_restantes} segundos"
    elif formato == "x min x sec":
        if horas > 0:
            return f"{horas} h {minutos} min {segundos_restantes} sec"
        else:
            return f"{minutos} min {segundos_restantes} sec"
    elif formato == "h:m:s":
        return f"{horas}:{minutos:02}:{segundos_restantes:02}"
    else:
        raise ValueError("Formato no reconocido. Usa uno de: 'x'x''', 'x minutos y x segundos', 'x min x sec', 'h:m:s'.")

# # Ejemplos de uso
# segundos_con_decimales = 3725.8  # 1 hora, 2 minutos y 5.8 segundos
#
# # Estilo x'x''
# print(convertir_segundos(segundos_con_decimales, "x'x''"))  # Salida: "1h 2'6''"
#
# # Estilo x minutos y x segundos
# print(convertir_segundos(segundos_con_decimales, "x minutos y x segundos"))  # Salida: "1 horas, 2 minutos y 6 segundos"
#
# # Estilo x min x sec
# print(convertir_segundos(segundos_con_decimales, "x min x sec"))  # Salida: "1 h 2 min 6 sec"
#
# # Estilo reloj h:m:s
# print(convertir_segundos(segundos_con_decimales, "h:m:s"))  # Salida: "1:02:06"


def separar_artistas(artistas):
    # Define regex pattern to capture the first match of " / ", " feat. ", or " canta: "
    pattern = re.compile(r" / | feat\. | canta: ")

    # Split the input string using the defined pattern, with maxsplit=1 to split only at the first occurrence
    artists = re.split(pattern, artistas, maxsplit=1)

    # If the second element exists but is empty, treat as if there's no split needed
    if len(artists) > 1 and artists[1].strip() == "":
        return artistas, ""

    # Assign artists based on the length of the split result
    artists1 = artists[0]  # Main artist
    artists2 = artists[1] if len(artists) > 1 else ""  # Second artist, if any

    return artists1, artists2


def unir_artistas(orquesta, cantor, caracter_separacion):
    # Verificamos si cantor es NaN o si es una cadena vacía o de espacios
    if cantor and not (isinstance(cantor, float) and math.isnan(cantor)) and cantor.strip():
        artista = f'{orquesta}{caracter_separacion}{cantor.strip()}'
    else:
        artista = orquesta

    return artista

def palabras_mas_comunes(db, columna):
    dataframe = db.copy()
    # Paso 2: Eliminar puntuaciones y dividir en palabras
    dataframe[columna] = dataframe[columna].apply(lambda x: re.findall(r'\b\w+\b', x))
    # Paso 3: Contar la frecuencia de las palabras
    word_count = Counter()
    dataframe[columna].apply(lambda x: word_count.update(x))
    # Obtener las palabras más comunes
    most_common_words = word_count.most_common()




def coincidencias_a_colores(bool_coincidencias, perfect_match):
    # Definir colores para cada tipo de coincidencia
    colores = {
        'EXACTO': 'lightgreen',
        'COINCIDENCIA': 'yellow',
        'PARCIAL': 'orange',
        'NEGATIVO': 'red'
    }

    # Crear un DataFrame vacío basado en los índices de los Series en bool_coincidencias
    indices = next(iter(bool_coincidencias.values())).index
    df_colores = pd.DataFrame(index=indices)

    # Si hay un perfect match, asignar todos los colores a gris y retornar
    if perfect_match:
        df_colores = pd.DataFrame('grey', index=indices, columns=[
            'ORQUESTA', 'TITULO', 'CANTOR', 'FECHA', 'ESTILO', 'COMPOSITOR'
        ])
        return df_colores

    # Mapear los colores para cada categoría relevante
    mapping = {
        'ORQUESTA': [TagLabels.ORQUESTA_EXACTA, TagLabels.ORQUESTA, TagLabels.ORQUESTA_PARCIAL,
                     TagLabels.ORQUESTA_NEGATIVO],
        'TITULO': [TagLabels.TITULO_EXACTO, TagLabels.TITULO, TagLabels.TITULO_PARCIAL, TagLabels.TITULO_NEGATIVO],
        'CANTOR': [TagLabels.CANTOR_EXACTO, TagLabels.CANTOR, TagLabels.CANTOR_PARCIAL, TagLabels.CANTOR_NEGATIVO],
        'FECHA': [TagLabels.FECHA_EXACTA, TagLabels.FECHA, TagLabels.FECHA_PARCIAL, TagLabels.FECHA_NEGATIVA],
        'ESTILO': [TagLabels.ESTILO_EXACTO, TagLabels.ESTILO, TagLabels.ESTILO_PARCIAL, TagLabels.ESTILO_NEGATIVO],
        'COMPOSITOR': [TagLabels.COMPOSITOR_AUTOR_EXACTO, TagLabels.COMPOSITOR_AUTOR,
                       TagLabels.COMPOSITOR_AUTOR_PARCIAL, TagLabels.COMPOSITOR_AUTOR_NEGATIVO]
    }

    # Procesar cada grupo de coincidencias
    for categoria, etiquetas in mapping.items():
        exacta, normal, parcial, negativo = etiquetas

        # Inicializar con el color de "Ninguna Coincidencia"
        df_colores[categoria] = colores['NEGATIVO']  # Asignar por defecto

        # Aplicar colores basados en las coincidencias
        for i in indices:
            if bool_coincidencias[exacta][i]:
                df_colores.at[i, categoria] = colores['EXACTO']
            elif bool_coincidencias[normal][i]:
                df_colores.at[i, categoria] = colores['COINCIDENCIA']
            elif bool_coincidencias[parcial][i]:
                df_colores.at[i, categoria] = colores['PARCIAL']
            elif bool_coincidencias[negativo][i]:
                df_colores.at[i, categoria] = colores['NEGATIVO']

    return df_colores


def update_tags(file_path, title=None, artist=None, year=None, genre=None, composer=None):
    # Get the file extension
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    if ext == '.mp3':
        # Load the MP3 file
        audio = MP3(file_path, ID3=EasyID3)
    elif ext == '.flac':
        # Load the FLAC file
        audio = FLAC(file_path)
    elif ext == '.m4a':
        # Load the M4A file
        audio = MP4(file_path)
    else:
        print(f"Unsupported file format: {ext}")
        return False

    # Update tags
    if title:
        if ext == '.m4a':
            audio['\xa9nam'] = title  # Tag for composer in M4A files
        else:
            audio['title'] = title
    if artist:
        if ext == '.m4a':
            audio['\xa9ART'] = artist  # Tag for composer in M4A files
        else:
            audio['artist'] = artist
    if year:
        if ext == '.m4a':
            audio['\xa9day'] = year  # Tag for composer in M4A files
        else:
            audio['date'] = year
    if genre:
        if ext == '.m4a':
            audio['\xa9gen'] = genre  # Tag for composer in M4A files
        else:
            audio['genre'] = genre
    if composer:
        if ext == '.m4a':
            audio['\xa9wrt'] = composer  # Tag for composer in M4A files
        else:
            audio['composer'] = composer

    # Save changes
    audio.save()
    return True


def capitalize_uppercase_words(text):
    words = text.split()
    transformed_words = []
    for word in words:

        if word == "DE" or word == "De":
            word = "de"
        if word == "DEL" or word == "Del":
            word = "del"
        if word == "DI" or word == "Di":
            word = "di"
        if word.isupper():
            word = word.capitalize()
            chars = list(word)
            for i in range(len(chars) - 1):
                # Check if the current character is a quote or hyphen
                if chars[i] in ["'", "-"]:
                    # Capitalize the next character
                    chars[i + 1] = chars[i + 1].upper()
            word = ''.join(chars)
            transformed_words.append(word)
        else:
            transformed_words.append(word)

    return ' '.join(transformed_words)


def aplicartag_archivo(ruta_archivo, coincidencias, coincidencia_preferida,coincidencia_titulo , tags):
    # Separar artistas originales en dos partes
    artists1, artists2 = separar_artistas(tags.artist)

    # Combinar artista y cantor preferidos
    artist = unir_artistas(
        coincidencias.artista.iloc[coincidencia_preferida],
        coincidencias.cantor.iloc[coincidencia_preferida],
        " / "
    )

    # Actualizar tags usando un solo llamado a la función update_tags
    update_tags(
        ruta_archivo,
        title=coincidencias.titulo.iloc[coincidencia_preferida],
        artist=artist,
        year=coincidencias.fecha.iloc[coincidencia_preferida],
        genre=coincidencias.estilo.iloc[coincidencia_preferida],
        composer=coincidencias.compositor_autor.iloc[coincidencia_preferida]
    )

    reporte_data = {
        'title': coincidencias.titulo.iloc[coincidencia_preferida],
        'artist': artist,
        'year': coincidencias.fecha.iloc[coincidencia_preferida],
        'genre': coincidencias.estilo.iloc[coincidencia_preferida],
        'composer': coincidencias.compositor_autor.iloc[coincidencia_preferida],
        'file_path': ruta_archivo,
        "Artista encontrado": False,
        "Titulo encontrado": False,
        "Numero de coincidencias": 1,
        "Hay coincidencia preferida": False,
        "No hay coincidencia preferida": False,
        "Coincidencia perfecta": True
    }

    # Crear el diccionario de reemplazo de tags
    reemplazo_tags_linea = crear_reemplazo_tags_linea(
        ruta_archivo=ruta_archivo,
        tags=tags,
        coincidencias=coincidencias,
        coincidencia_preferida=coincidencia_preferida,
        artists1=artists1,
        artists2=artists2,
        coincidencia_titulo=coincidencia_titulo
    )

    return reemplazo_tags_linea, reporte_data



def crear_reemplazo_tags_linea(
    ruta_archivo, tags, coincidencias, coincidencia_preferida, artists1, artists2, coincidencia_titulo
):
    """
    Crea un diccionario de reemplazo de tags para una coincidencia preferida.

    :param ruta_archivo: Ruta del archivo.
    :param tags: Objeto con los tags actuales (title, artist, year, genre, composer).
    :param coincidencias: DataFrame con las coincidencias.
    :param coincidencia_preferida: Índice de la coincidencia preferida.
    :param artists1: Artista original.
    :param artists2: Cantor original.
    :return: Diccionario con los reemplazos de tags.
    """
    if coincidencia_preferida is not None:
        data = {
            'archivo': [ruta_archivo],
            'old_title': [tags.title],
            'new_title': [coincidencias.titulo.iloc[coincidencia_preferida]],
            'old_artist': [artists1],
            'new_artist': [coincidencias.artista.iloc[coincidencia_preferida]],
            'old_cantor': [artists2],
            'new_cantor': [coincidencias.cantor.iloc[coincidencia_preferida]],
            'old_year': [tags.year],
            'new_year': [coincidencias.fecha.iloc[coincidencia_preferida]],
            'old_genre': [tags.genre],
            'new_genre': [coincidencias.estilo.iloc[coincidencia_preferida]],
            'old_composer': [tags.composer],
            'new_composer': [coincidencias.compositor_autor.iloc[coincidencia_preferida]],
            'coincidencia_titulo': [coincidencia_titulo]
        }
        return pd.DataFrame(data)
    else:
        return pd.DataFrame()  # Return an empty DataFrame if no preferred match


def link_to_music(link):
    path = link.split('/', 3)[-1]
    local_path = os.path.join(mp3_dir, path)
    return local_path


def buscar_preferencias(booleanos, buscar_fecha):
    get = lambda label, idx=0: booleanos[label].iloc[idx]

    mapping = {
        'ORQUESTA': [TagLabels.ORQUESTA_EXACTA, TagLabels.ORQUESTA, TagLabels.ORQUESTA_PARCIAL,
                     TagLabels.ORQUESTA_NEGATIVO],
        'TITULO': [TagLabels.TITULO_EXACTO, TagLabels.TITULO, TagLabels.TITULO_PARCIAL, TagLabels.TITULO_NEGATIVO],
        'CANTOR': [TagLabels.CANTOR_EXACTO, TagLabels.CANTOR, TagLabels.CANTOR_PARCIAL, TagLabels.CANTOR_NEGATIVO],
        'FECHA': [TagLabels.FECHA_EXACTA, TagLabels.FECHA, TagLabels.FECHA_PARCIAL, TagLabels.FECHA_NEGATIVA],
        'ESTILO': [TagLabels.ESTILO_EXACTO, TagLabels.ESTILO, TagLabels.ESTILO_PARCIAL, TagLabels.ESTILO_NEGATIVO],
        'COMPOSITOR': [TagLabels.COMPOSITOR_AUTOR_EXACTO, TagLabels.COMPOSITOR_AUTOR,
                       TagLabels.COMPOSITOR_AUTOR_PARCIAL, TagLabels.COMPOSITOR_AUTOR_NEGATIVO]
    }
    # Procesar cada grupo de coincidencias
    for categoria, etiquetas in mapping.items():
        exacta, normal, parcial, negativo = etiquetas

    longitudes = [len(series) for series in booleanos.values()]
         # Verifica si todas las longitudes son iguales
    longitudes_iguales = len(set(longitudes)) == 1
         # Retorna la longitud común si son iguales, o None si no lo son
    if longitudes_iguales:
        numero_de_resultados = longitudes[0]  # Todas las longitudes son iguales, devuelve una de ellas
    else:
        print('ERROR, las lista de booleanos no son iguales')
        return False, 0

    # Corrección y simplificación del if
    if (numero_de_resultados == 1 and
            (get(TagLabels.TITULO_EXACTO, 0) or get(TagLabels.TITULO, 0)) and
            (get(TagLabels.ORQUESTA_EXACTA, 0) or get(TagLabels.ORQUESTA, 0)) and
            not get(TagLabels.FECHA_NEGATIVA, 0) and not not get(TagLabels.CANTOR_NEGATIVO, 0)):
        return True, 0



        # Comprobar las condiciones para múltiples resultados
    preferidos = []
    for idx in range(numero_de_resultados):
        titulo_condicion = get(TagLabels.TITULO_EXACTO, idx) or get(TagLabels.TITULO, idx)
        orquesta_condicion = get(TagLabels.ORQUESTA_EXACTA, idx) or get(TagLabels.ORQUESTA, idx)
        fecha_condicion = not get(TagLabels.FECHA_NEGATIVA, idx) if buscar_fecha else True
        cantor_condicion = not get(TagLabels.CANTOR_NEGATIVO, idx)

        if titulo_condicion and orquesta_condicion and fecha_condicion and cantor_condicion:
            preferidos.append(idx)

    # Si hay exactamente un preferido, devolverlo
    if len(preferidos) == 1:
        return True, preferidos[0]


    return False, 0


def es_par(numero):
    """Devuelve True si el número es par, de lo contrario False."""
    return numero % 2 == 0


def convert_numbers_to_words(text):
    # Find all numbers in the text
    numbers = re.findall(r'\d+', text)
    # Replace each number with its word form
    for number in numbers:
        text = text.replace(number, nw.num2words(number, lang='es'))

    return text


def contain_most_words_in_dic(dictionary, text):
    # Limpia y prepara el texto
    text = unidecode(convert_numbers_to_words(text)).lower().strip()
    text = re.sub(r"[()]", "", text)  # Elimina paréntesis
    text_words = set(text.split())

    lista_numero_palabras_comun = []

    for key in dictionary:
        words = key
        words = re.sub(r"[()]", "", words)  # Limpia la clave
        words_words = set(words.lower().split())

        # Encuentra las palabras comunes entre el texto y la clave del diccionario
        common_words = text_words.intersection(words_words)
        numero_de_palabras_en_comun = len(common_words)

        # Agrega la clave y el número de palabras comunes
        lista_numero_palabras_comun.append((key, numero_de_palabras_en_comun))

    # Encuentra el máximo número de palabras en común
    maximo_palabras = max(lista_numero_palabras_comun, key=lambda x: x[1])[1]

    # Encuentra las claves que tienen ese máximo número de palabras en común
    if maximo_palabras > 1:
        keys_mas_palabras = [key for key, value in lista_numero_palabras_comun if value == maximo_palabras]
        return keys_mas_palabras
    else:
        return None

    # Encuentra los índices que tienen ese máximo número de palabras en común
    indices_mas_palabras = [index for index, value in lista_numero_palabras_comun if value == maximo_palabras]

    # Si hay exactamente una coincidencia, devuelve la cadena correspondiente
    if len(indices_mas_palabras) == 1 and maximo_palabras > 0:
        return list(dictionary.keys())[indices_mas_palabras[0]]
    else:
        return None  # Si hay más de una coincidencia o ninguna, no devuelve nada


def contain_most_words(database, text, columna):
    # Normalize and preprocess the input text
    text = unidecode(convert_numbers_to_words(text)).lower().strip()
    text = text.replace("(", "").replace(")", "")
    text_words = set(word for word in text.lower().split() if word not in articulos_preposiciones_comunes)

    # Lists to store counts of common words
    lista_numero_palabras_comun = []

    for index, row in database.iterrows():
        # Normalize and preprocess the row text
        words = row[columna]
        words = words.replace("(", "").replace(")", "")
        words_words = set(words.lower().split())

        # Calculate the number of common words
        common_words = text_words.intersection(words_words)
        numero_de_palabras_en_comun = len(common_words)
        lista_numero_palabras_comun.append((index, numero_de_palabras_en_comun))

    # Find the maximum number of common words
    maximo_palabras = max(lista_numero_palabras_comun, key=lambda x: x[1])[1]

    # Create a boolean Series aligned with the DataFrame's index
    coincidencias = pd.Series(False, index=database.index)
    if maximo_palabras > 0:
        # Find indices with the maximum number of common words
        indices_mas_palabras = [index for index, value in lista_numero_palabras_comun if value == maximo_palabras]
        coincidencias.loc[indices_mas_palabras] = True

    return coincidencias

def get_file_name_without_extension(file_path):
    # Extract the file name with extension
    file_name_with_ext = os.path.basename(file_path)
    # Split the file name and extension, and return the file name
    file_name, _ = os.path.splitext(file_name_with_ext)
    return file_name


def buscar_titulo(database, tag):
    """
    Busca un título en una base de datos utilizando diversas estrategias de coincidencia, con optimizaciones de rendimiento.

    Proceso paso a paso:
    1. Normaliza y limpia el título original a buscar.
    2. Realiza coincidencias exactas y normalizadas:
       - Coincidencia exacta con el título original.
       - Coincidencia exacta con el título normalizado.
    3. Busca coincidencias parciales y patrones en los títulos:
       - Coincidencias al inicio de los títulos.
       - Coincidencias en cualquier parte de los títulos.
       - Coincidencias donde los títulos empiezan con el término buscado.
    4. Coincidencia basada en la mayor cantidad de palabras coincidentes.
    5. Devuelve el tipo de coincidencia y las entradas correspondientes.

    :param database: DataFrame que contiene los títulos donde buscar.
    :param tag: Objeto con el título a buscar y metadatos adicionales.
    :return: Tupla que contiene un indicador del tipo de coincidencia y un DataFrame con las entradas coincidentes.
    """

    # Normaliza y limpia el título original; si está vacío, usa el nombre del archivo sin extensión
    titulo_original = ftfy.fix_text(tag.title).strip() if tag.title else get_file_name_without_extension(tag._filename)
    titulo_buscar_min = unidecode(convert_numbers_to_words(titulo_original)).lower().strip()

    # Verifica si el título es válido para la búsqueda
    if titulo_buscar_min == "":
        return 0, pd.DataFrame()  # No hay título válido para buscar

    # Coincidencias exactas

    # Paso 2.1: Coincidencia exacta del título original
    # Variable clave: coincidencias (boolean series que compara los títulos originales)
    coincidencias = database["titulo"] == titulo_original
    if coincidencias.any():
        return 6, database[coincidencias].copy()  # Retorna 6 si hay coincidencias exactas con el título original

    # Paso 2.2: Coincidencia exacta con el título normalizado
    # Variable clave: coincidencias (boolean series que compara los títulos normalizados)
    coincidencias = database["titulo_min"] == titulo_buscar_min
    if coincidencias.any():
        return 5, database[coincidencias].copy()  # Retorna 5 si hay coincidencias exactas con el título normalizado

    # Coincidencias parciales y patrones en los títulos

    # Paso 3.1: Coincidencias al inicio de los títulos
    # Variable clave: coincidencias (boolean series que verifica si los títulos comienzan con el término de búsqueda usando regex)
    coincidencias = database["titulo_min"].str.contains(f'^{titulo_buscar_min}', case=False, na=False, regex=False)
    if coincidencias.any():
        return 4, database[
            coincidencias].copy()  # Retorna 4 si el término de búsqueda coincide al inicio de algún título


    # Paso 3.3: Coincidencias donde los títulos empiezan con el término buscado
    # Variable clave: coincidencias (boolean series que verifica si los títulos comienzan exactamente con el término de búsqueda)
    coincidencias = database["titulo_min"].str.startswith(titulo_buscar_min, na=False)
    if coincidencias.any():
        # Encuentra la coincidencia más larga por la longitud del título
        indice_mejor_coincidencia = database.loc[coincidencias, "titulo_min"].str.len().idxmax()
        coincidencias = database.index == indice_mejor_coincidencia

        # Guarda información residual si está configurado
        if guardar_residuos:
            titulo_coincidencia = database.loc[coincidencias, 'titulo_min'].iloc[0]
            nuevos_residuos = [tag.title, titulo_coincidencia, tag.title[len(titulo_coincidencia):], tag._filename]
            residuos.append(nuevos_residuos)

        return 3, database[
            coincidencias].copy()  # Retorna 2 si encuentra el título más largo que empieza con el término de búsqueda

    # Paso 3.2: Coincidencias en cualquier parte de los títulos
    # Variable clave: coincidencias (boolean series que verifica si los títulos contienen el término de búsqueda en cualquier parte)
    coincidencias = database["titulo_min"].str.contains(titulo_buscar_min, case=False, na=False, regex=False)
    if coincidencias.any():
        return 2, database[
            coincidencias].copy()  # Retorna 3 si el término de búsqueda se encuentra en cualquier parte del título


    # Coincidencia basada en palabras

    # Paso 4.1: Coincidencia basada en la mayor cantidad de palabras coincidentes
    # Variable clave: coincidencias (series booleana generada por la función personalizada `contain_most_words`)
    coincidencias = pd.Series(contain_most_words(database, titulo_buscar_min, "titulo_min"), index=database.index)
    if coincidencias.any():
        return 1, database[
            coincidencias].copy()  # Retorna 1 si hay coincidencias basadas en la mayor cantidad de palabras coincidentes

    # Si no se encontraron coincidencias
    return 0, pd.DataFrame()  # Retorna 0 si no se encontraron coincidencias


def compare_tags(artista_coincidencia, titulo_coincidencia, database, tag):
    coincidencias = {}
    artista_original, cantor_original = separar_artistas(tag.artist)
    artista_buscar = unidecode(artista_original).lower().strip()
    cantor_buscar = unidecode(cantor_original).lower()

    database_length = len(database)

    # Coincidencias para orquesta
    coincidencias[TagLabels.ORQUESTA_EXACTA] = (database["artista"] == artista_original)
    coincidencias[TagLabels.ORQUESTA] = pd.Series([artista_coincidencia == 2] * database_length, index=database.index)
    coincidencias[TagLabels.ORQUESTA_PARCIAL] = pd.Series([artista_coincidencia == 1] * database_length, index=database.index)
    coincidencias[TagLabels.ORQUESTA_NEGATIVO] = pd.Series([artista_coincidencia == 0] * database_length, index=database.index)

    # Coincidencias para título
    coincidencias[TagLabels.TITULO_EXACTO] = pd.Series([titulo_coincidencia == 6] * database_length, index=database.index)
    coincidencias[TagLabels.TITULO] = pd.Series([titulo_coincidencia in [2,3,4,5]] * database_length, index=database.index)
    coincidencias[TagLabels.TITULO_PARCIAL] = pd.Series([titulo_coincidencia in [1]] * database_length,
                                                        index=database.index)  # Coincidencia parcial
    coincidencias[TagLabels.TITULO_NEGATIVO] = pd.Series([titulo_coincidencia == 0] * database_length, index=database.index)

    # Coincidencias para cantor
    coincidencias[TagLabels.CANTOR_EXACTO] = (database["cantor"] == cantor_original)
    coincidencias[TagLabels.CANTOR] = database["cantor_min"].str.contains(cantor_buscar, case=False, na=False, regex=False)
    coincidencias[TagLabels.CANTOR_PARCIAL] = pd.Series([cantor_buscar == ""] * database_length, index=database.index)
    coincidencias[TagLabels.CANTOR_NEGATIVO] = pd.Series([True] * database_length, index=database.index)

    # Coincidencias para fecha
    tag.year = tag.year or ""
    tag.genre = tag.genre or ""
    tag.composer = tag.composer or ""

    # Extraer el año
    if extraer_cuatro_numeros(tag.year):
        ano = int(extraer_cuatro_numeros(tag.year))
    else:
        ano = ""

    # Comprobación de coincidencias exactas
    coincidencias[TagLabels.FECHA_EXACTA] = database["fecha"] == tag.year

    coincidencias[TagLabels.FECHA] = database["fecha_ano"] == ano

    # Intentar obtener la fecha de fallecimiento del director
    try:
        fallecimiento_director = fallecimiento[artista_buscar]
    except KeyError:
        # Si el artista no está en la lista, asignamos la fecha actual como indicación de "futuro"
        fallecimiento_director = datetime.now().date()
        # print(f"Fallecimiento de {artista_buscar} NO encontrado: {fallecimiento_director}")


    # Intentamos convertir tag.year a una fecha
    tag_year_date = parse_year(tag.year)

    # Si la fecha no es None, realizamos la comparación; si es None, asumimos que no falleció
    if tag_year_date is not None:
        Fallecio = fallecimiento_director < tag_year_date
    else:
        Fallecio = False  # Si no se puede obtener la fecha, no se considera que falleció

    if not tag.year or tag.year == "" or Fallecio:
        coincidencias[TagLabels.FECHA_PARCIAL] = pd.Series([True] * database_length, index=database.index)
    else:
        coincidencias[TagLabels.FECHA_PARCIAL] = pd.Series([False] * database_length, index=database.index)
    coincidencias[TagLabels.FECHA_NEGATIVA] = database["fecha_ano"] != extraer_cuatro_numeros(tag.year)

    # Coincidencias para estilo
    coincidencias[TagLabels.ESTILO_EXACTO] = database["estilo"] == tag.genre
    coincidencias[TagLabels.ESTILO] = pd.Series([False] * database_length, index=database.index)
    coincidencias[TagLabels.ESTILO_PARCIAL] = pd.Series([False] * database_length, index=database.index)
    coincidencias[TagLabels.ESTILO_NEGATIVO] = pd.Series([True] * database_length, index=database.index)

    # Coincidencias para compositor_autor
    coincidencias[TagLabels.COMPOSITOR_AUTOR_EXACTO] = database["compositor_autor"] == tag.composer
    coincidencias[TagLabels.COMPOSITOR_AUTOR] = pd.Series([False] * database_length, index=database.index)
    coincidencias[TagLabels.COMPOSITOR_AUTOR_PARCIAL] = pd.Series([False] * database_length, index=database.index)
    coincidencias[TagLabels.COMPOSITOR_AUTOR_NEGATIVO] = pd.Series([True] * database_length, index=database.index)

    #Ajustes para asegurar que solo uno sea True por cada grupo
    def ajustar_grupo(grupo):
        exacto, completo, parcial, negativo = grupo
        coincidencias[completo] = coincidencias[completo] & ~coincidencias[exacto]
        coincidencias[parcial] = coincidencias[parcial] & ~(coincidencias[exacto] | coincidencias[completo])
        coincidencias[negativo] = coincidencias[negativo] & ~(
                    coincidencias[exacto] | coincidencias[completo] | coincidencias[parcial])

    # Ajustar todos los grupos
    ajustar_grupo(
        [TagLabels.ORQUESTA_EXACTA, TagLabels.ORQUESTA, TagLabels.ORQUESTA_PARCIAL, TagLabels.ORQUESTA_NEGATIVO])
    ajustar_grupo([TagLabels.TITULO_EXACTO, TagLabels.TITULO, TagLabels.TITULO_PARCIAL, TagLabels.TITULO_NEGATIVO])
    ajustar_grupo([TagLabels.CANTOR_EXACTO, TagLabels.CANTOR, TagLabels.CANTOR_PARCIAL, TagLabels.CANTOR_NEGATIVO])
    ajustar_grupo([TagLabels.FECHA_EXACTA, TagLabels.FECHA, TagLabels.FECHA_PARCIAL, TagLabels.FECHA_NEGATIVA])
    ajustar_grupo([TagLabels.ESTILO_EXACTO, TagLabels.ESTILO, TagLabels.ESTILO_PARCIAL, TagLabels.ESTILO_NEGATIVO])
    ajustar_grupo([TagLabels.COMPOSITOR_AUTOR_EXACTO, TagLabels.COMPOSITOR_AUTOR, TagLabels.COMPOSITOR_AUTOR_PARCIAL,
                   TagLabels.COMPOSITOR_AUTOR_NEGATIVO])

    # Check all fields (todo)
    coincidencias[TagLabels.TODO] = (
            coincidencias[TagLabels.TITULO_EXACTO] &
            coincidencias[TagLabels.ORQUESTA_EXACTA] &
            coincidencias[TagLabels.CANTOR_EXACTO] &
            coincidencias[TagLabels.FECHA_EXACTA] &
            coincidencias[TagLabels.ESTILO_EXACTO] &
            coincidencias[TagLabels.COMPOSITOR_AUTOR_EXACTO]
    )

    # Determinar si hay un perfect match y encontrar el índice de la coincidencia
    perfect_match = coincidencias[TagLabels.TODO].any()
    perfect_match_index = coincidencias[TagLabels.TODO].idxmax() if perfect_match else None

    return coincidencias, perfect_match, perfect_match_index



def stop_music():
    pygame.mixer.music.stop()


def concaternar_autores(compositor, autor):
    # if not isinstance(compositor, str):
    #     compositor = "?"
    # if not isinstance(autor, str):
    #     autor = "?"
    if compositor != "" and autor != "":
        concatenacion = "Musica: " + compositor + " - Letra: " + autor
    elif compositor != "":
        concatenacion = "Musica: " + compositor
    elif autor != "":
        concatenacion = "Letra: " + autor
    else:
        concatenacion = "Desconocido"

    return concatenacion


def adjust_text_size(text_frame, max_width_cm, max_font_size=100, min_font_size=10):
    # Convertir las medidas de cm a puntos (pt)
    max_width_pt = max_width_cm / 11500

    font_size = max_font_size
    text = ''.join(run.text for paragraph in text_frame.paragraphs for run in paragraph.runs)
    text_length = len(text)

    while font_size >= min_font_size:
        # Aplicar el tamaño de la fuente
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                run.font.name = DEFAULT_FONT_NAME

        # Cálculo del ancho estimado del texto basado en el tamaño de la fuente y la longitud del texto
        char_width = (DEFAULT_CHAR_WIDTH / 100) * font_size
        estimated_text_width = char_width * text_length

        # Si el ancho estimado del texto es menor o igual al ancho máximo permitido, se rompe el bucle
        if estimated_text_width <= max_width_pt:
            break

        font_size -= 1

    # Aplicar el tamaño de fuente final
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.name = DEFAULT_FONT_NAME


def obtener_intervalo_anos(lista):
    # Filtramos sublistas vacías y convertimos los años a enteros, ignorando strings vacíos
    años_int = [int(sublista[0]) for sublista in lista if sublista and sublista[0].isdigit()]
    # Si la lista resultante está vacía, retornamos un mensaje de error o un valor predeterminado
    if not años_int:
        return "No hay años válidos en la lista"

    # Encontramos el año mínimo y máximo
    año_min = min(años_int)
    año_max = max(años_int)
    # Retornamos el string en el formato deseado
    return f"({año_min} - {año_max})"



def obtener_autores(lista):
    # Convertimos la lista en un conjunto para eliminar duplicados
    autores = set()

    for sublista in lista:
        for item in sublista:
            if isinstance(item, str) and item != 'nan':
                # Dividimos el string por ' y ' y ', ' para manejar los casos de varios autores
                partes = item.split(' y ')
                for parte in partes:
                    sub_partes = parte.split(', ')
                    for sub_parte in sub_partes:
                        autores.add(sub_parte.strip())

    # Si no hay autores, devolvemos 'instrumental'
    if not autores:
        return 'instrumental'

    # Convertimos el conjunto a lista para ordenarlos de manera consistente
    autores = list(autores)

    # Ordenamos para mantener consistencia
    autores.sort()

    # Según la cantidad de autores, retornamos la cadena deseada
    if len(autores) == 1:
        return autores[0]
    elif len(autores) == 2:
        return ' y '.join(autores)
    else:
        return ', '.join(autores[:-1]) + ' y ' + autores[-1]

def get_largest_paragraph(text):
    """
    Separates the input text into paragraphs based on newline characters and returns the largest paragraph.

    Parameters:
        text (str): The input string containing paragraphs separated by newline characters.

    Returns:
        str: The largest paragraph based on the length of characters.
    """
    # Split the text into paragraphs using newline as the separator
    paragraphs = text.split('\n')

    # Remove empty paragraphs (in case of multiple consecutive newline characters)
    paragraphs = [para.strip() for para in paragraphs if para.strip()]

    # Find the largest paragraph based on length
    largest_paragraph = max(paragraphs, key=len) if paragraphs else ""

    return largest_paragraph


def guardar_archivo_output(tipo, dataframe, encabezados=None):
    now = datetime.now()
    timestamp = now.strftime('%Y%m%d_%H%M%S')
    filename = f'{tipo}_{timestamp}.csv'
    file_path = os.path.join(output_folder, filename)

    # Open the file to write data
    with open(file_path, 'w', newline='', encoding='utf-8') as file:
        # If custom headers are provided, write them manually
        if encabezados is not None:
            file.write(encabezados + '\n')  # Write the custom header line with the origin
            # Write the DataFrame without the default headers since they are custom
            dataframe.to_csv(file, index=False, sep=';', header=True)
        else:
            # Write the DataFrame with its own headers
            dataframe.to_csv(file, index=False, sep=';', header=True)

    print(f'{tipo.capitalize()} guardados en: {file_path}')

    return file_path


def leer_tags(ruta_archivo):
    tags = TinyTag.get(ruta_archivo)
    return tags

def find_similar_lines(merged_lines, similarity_threshold=0.97):
    """
    Find and print lines in merged_lines that are very similar to each other.
    Lines are considered similar if their similarity ratio is above the given threshold.
    """
    similar_pairs = []

    # Compare each line with every other line
    for i, line1 in enumerate(merged_lines):
        for j, line2 in enumerate(merged_lines):
            if i != j:  # Avoid comparing the same line with itself
                # Calculate similarity ratio
                similarity = difflib.SequenceMatcher(None, line1, line2).ratio()
                if similarity >= similarity_threshold:
                    # Add the pair if they are similar enough
                    similar_pairs.append((i, line1, j, line2, similarity))

    if similar_pairs:
        print("Similar lines found:")
        for pair in similar_pairs:
            print(f"Line {pair[0]}: {pair[1]}")
            print(f"Line {pair[2]}: {pair[3]}")
            print(f"Similarity: {pair[4]:.2f}\n")
    else:
        print("No similar lines found.")


from datetime import datetime

def parse_year(tag_year):
    # Definir los formatos de fecha que se van a intentar
    formatos_fecha = ["%Y-%m-%d", "%d/%m/%Y", "%d-%m-%Y", "%Y", "%Y%m%d"]

    # Si el valor es un entero, lo tratamos como un año
    if isinstance(tag_year, int):
        return datetime(tag_year, 1, 1).date()

    # Si es una cadena, intentamos convertirla a una fecha usando múltiples formatos
    elif isinstance(tag_year, str):
        tag_year = tag_year.strip()  # Limpiamos espacios

        # Intentamos cada formato de fecha
        for formato in formatos_fecha:
            try:
                parsed_date = datetime.strptime(tag_year, formato).date()
                return parsed_date
            except ValueError:
                # Si falla, seguimos intentando con otros formatos
                continue

    # Si no se pudo convertir a fecha válida, retornamos None
    return None



def extraer_cuatro_numeros(cadena):
    # Usamos la expresión regular \d{4} para encontrar cuatro dígitos consecutivos
    resultado = re.search(r'\d{4}', cadena)
    if resultado:
        return resultado.group()
    else:
        return None

def generar_timestamp():
    return datetime.now().strftime("%Y%m%d%H%M%S")


def extract_year(date_str):
    # Define a regular expression pattern to match the year
    pattern = r'\b(\d{4})\b'

    # Match the pattern
    match = re.search(pattern, date_str)

    # If there's a match, return the year
    if match:
        return match.group(1)
    else:
        return ""


def convert_date_format(date_str):
    """
    Convert a date from DD/MM/YYYY format to YYYY-MM-DD format.
    Args:
        date_str (str): Date string in DD/MM/YYYY format.
    Returns:
        str: Date string in YYYY-MM-DD format.
    """
    pattern_yyyymmdd = r'^[0-9]{4}-(0[0-9]|1[0-2])-(0[0-9]|[12][0-9]|3[01])$'

    if re.match(pattern_yyyymmdd, date_str):
        if date_str.endswith('-00-00'):
            # Remove the last three characters
            return date_str[:-6]
        if date_str.endswith('-00'):
            # Remove the last three characters
            return date_str[:-3]
        new_date_str = date_str
    else:
        # Parse the input date string
        date_object = datetime.strptime(date_str, "%d/%m/%Y")

        # Convert to the desired format
        new_date_str = date_object.strftime("%Y-%m-%d")

    return new_date_str

