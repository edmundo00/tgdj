import re
from unidecode import unidecode
from config import *
from pptx.util import Pt, Cm


def separar_artistas(artistas):
    artists = artistas.split(" / ")
    artists1 = artists[0]
    if len(artists) > 1:
        artists2 = artists[1]
    else:
        artists2 = ""
    return artists1, artists2


def palabras_mas_comunes(db, columna):
    dataframe = db.copy()
    # Paso 2: Eliminar puntuaciones y dividir en palabras
    dataframe[columna] = dataframe[columna].apply(lambda x: re.findall(r'\b\w+\b', x))
    # Paso 3: Contar la frecuencia de las palabras
    word_count = Counter()
    dataframe[columna].apply(lambda x: word_count.update(x))
    # Obtener las palabras más comunes
    most_common_words = word_count.most_common()


def extract_year(date_str):
    # Define a regular expression pattern to match the year
    pattern = r'\b(\d{4})\b'

    # Match the pattern
    match = re.search(pattern, date_str)

    # If there's a match, return the year
    if match:
        return match.group(1)
    else:
        return ""


def convert_date_format(date_str):
    """
    Convert a date from DD/MM/YYYY format to YYYY-MM-DD format.
    Args:
        date_str (str): Date string in DD/MM/YYYY format.
    Returns:
        str: Date string in YYYY-MM-DD format.
    """
    pattern_yyyymmdd = r'^[0-9]{4}-(0[0-9]|1[0-2])-(0[0-9]|[12][0-9]|3[01])$'

    if re.match(pattern_yyyymmdd, date_str):
        if date_str.endswith('-00-00'):
            # Remove the last three characters
            return date_str[:-6]
        if date_str.endswith('-00'):
            # Remove the last three characters
            return date_str[:-3]
        new_date_str = date_str
    else:
        # Parse the input date string
        date_object = datetime.strptime(date_str, "%d/%m/%Y")

        # Convert to the desired format
        new_date_str = date_object.strftime("%Y-%m-%d")

    return new_date_str


def update_tags(file_path, title=None, artist=None, year=None, genre=None, composer=None):
    # Get the file extension
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    if ext == '.mp3':
        # Load the MP3 file
        audio = MP3(file_path, ID3=EasyID3)
    elif ext == '.flac':
        # Load the FLAC file
        audio = FLAC(file_path)
    elif ext == '.m4a':
        # Load the M4A file
        audio = MP4(file_path)
    else:
        print(f"Unsupported file format: {ext}")
        return

    # Update tags
    if title:
        if ext == '.m4a':
            audio['\xa9nam'] = title  # Tag for composer in M4A files
        else:
            audio['title'] = title
    if artist:
        if ext == '.m4a':
            audio['\xa9ART'] = artist  # Tag for composer in M4A files
        else:
            audio['artist'] = artist
    if year:
        if ext == '.m4a':
            audio['\xa9day'] = year  # Tag for composer in M4A files
        else:
            audio['date'] = year
    if genre:
        if ext == '.m4a':
            audio['\xa9gen'] = genre  # Tag for composer in M4A files
        else:
            audio['genre'] = genre
    if composer:
        if ext == '.m4a':
            audio['\xa9wrt'] = composer  # Tag for composer in M4A files
        else:
            audio['composer'] = composer

    # Save changes
    audio.save()


def capitalize_uppercase_words(text):
    words = text.split()
    transformed_words = []
    for word in words:

        if word == "DE" or word == "De":
            word = "de"
        if word == "DEL" or word == "Del":
            word = "del"
        if word == "DI" or word == "Di":
            word = "di"
        if word.isupper():
            word = word.capitalize()
            chars = list(word)
            for i in range(len(chars) - 1):
                # Check if the current character is a quote or hyphen
                if chars[i] in ["'", "-"]:
                    # Capitalize the next character
                    chars[i + 1] = chars[i + 1].upper()
            word = ''.join(chars)
            transformed_words.append(word)
        else:
            transformed_words.append(word)

    return ' '.join(transformed_words)


def extraer_cuatro_numeros(cadena):
    # Usamos la expresión regular \d{4} para encontrar cuatro dígitos consecutivos
    resultado = re.search(r'\d{4}', cadena)
    if resultado:
        return resultado.group()
    else:
        return None


def link_to_music(link):
    path = link.split('/', 3)[-1]
    local_path = os.path.join(mp3_dir, path)
    return local_path


def es_par(numero):
    """Devuelve True si el número es par, de lo contrario False."""
    return numero % 2 == 0


def convert_numbers_to_words(text):
    # Find all numbers in the text
    numbers = re.findall(r'\d+', text)
    # Replace each number with its word form
    for number in numbers:
        text = text.replace(number, nw.num2words(number, lang='es'))

    text = unidecode(text).lower()

    return text


def contain_most_words(database, text, columna):
    text = unidecode(convert_numbers_to_words(text)).lower()
    text = text.replace("(", "").replace(")", "")
    text_words = set(text.lower().split())
    lista_numero_palabras_comun = []
    words_filas = []

    for index, row in database.iterrows():
        words = row[columna]
        words = words.replace("(", "").replace(")", "")
        words_words = set(words.lower().split())
        words_filas.append([index, words_words])
        common_words = text_words.intersection(words_words)
        numero_de_palabras_en_comun = len(common_words)
        lista_numero_palabras_comun.append((index, numero_de_palabras_en_comun))

    # Encuentra el máximo número de palabras en común
    maximo_palabras = max(lista_numero_palabras_comun, key=lambda x: x[1])[1]

    # Encuentra los índices que tienen ese máximo número de palabras en común
    if maximo_palabras > 0:
        indices_mas_palabras = [index for index, value in lista_numero_palabras_comun if value == maximo_palabras]
    else:
        indices_mas_palabras = []

    return indices_mas_palabras


def coincide(database, tag, que_coincide, talcual=False):
    lista_de_booleanos = []

    artista_original, cantor_original = separar_artistas(tag.artist)

    tituloabuscar = unidecode(convert_numbers_to_words(tag.title)).lower().strip()
    artistaabuscar = unidecode(convert_numbers_to_words(artista_original)).lower()
    cantorabuscar = unidecode(convert_numbers_to_words(cantor_original)).lower()

    if que_coincide == 'titulo':
        valor_buscado = tag.title
        if talcual:
            lista_de_booleanos = database["titulo"] == tag.title
        else:
            lista_de_booleanos = database["titulo_min"].str.contains(tituloabuscar, case=False, na=False, regex=False)

    if que_coincide == 'artista':
        valor_buscado = artista_original
        if talcual:
            lista_de_booleanos = database["artista"] == artista_original
        else:
            lista_de_booleanos = database["artista_min"].str.contains(artistaabuscar, case=False, na=False,
                                                                      regex=False)
    if que_coincide == 'cantor':
        valor_buscado = cantor_original
        if talcual:
            lista_de_booleanos = database["cantor"] == cantor_original
        else:
            lista_de_booleanos = database["cantor_min"].str.contains(cantorabuscar, case=False, na=False,
                                                                     regex=False)

    if tag.year == None:
        tag.year = ""
    if que_coincide == 'fecha':
        valor_buscado = tag.year
        lista_de_booleanos = database["fecha"] == tag.year

    if que_coincide == 'ano':
        lista_de_booleanos = database["fecha_ano"] == extraer_cuatro_numeros(tag.year)
        valor_buscado = extraer_cuatro_numeros(tag.year)

    if que_coincide == 'genero':
        valor_buscado = tag.genre
        if talcual:
            lista_de_booleanos = database["estilo"] == tag.genre
        else:
            lista_de_booleanos = database.estilo.str.contains(artistaabuscar, case=False, na=False,
                                                              regex=False)

    if que_coincide == 'compositor_autor':
        valor_buscado = tag.composer
        lista_de_booleanos = database["compositor_autor"] == tag.composer

    if que_coincide == 'todo':
        ti = database["titulo"]
        ar = database["artista"]
        ca = database["cantor"]
        fe = database["fecha"]
        es = database["estilo"]
        co = database["compositor_autor"]

        lista_de_booleanos = (ti == tag.title) & (ar == artista_original) & (ca == cantor_original) & (
                    fe == tag.year) & (es == tag.genre) & (co == tag.composer)

    return lista_de_booleanos


def stop_music():
    pygame.mixer.music.stop()

def limpiar_base_dato():

    # db = pd.read_csv(csv_grabaciones, encoding="utf-8", sep=";")
    # quitar los valores ' (2)', ' (3)', ' (4)', ' (5)', ' (b)', ' (c)' del titulo
    # quitar_de_titulos = [' (2)', ' (3)', ' (4)', ' (5)', ' (b)', ' (c)']
    # for palabra in quitar_de_titulos:
    #     db['titulo'] = db['titulo'].str.replace(palabra, "", regex=False)
    #
    # # poner estilos Tango, Tango Milonga, Tango Vals.
    # replacements = [
    #     ('TANGO', 'tango'),
    #     ('VALS', 'tango vals'),
    #     ('MILONGA', 'tango milonga')
    # ]
    # replacement_dict = dict(replacements)
    # db['estilo'] = db['estilo'].replace(replacement_dict)
    # db['estilo'] = db['estilo'].apply(lambda x: x.lower())
    #
    # # reemplazar los nan con cadenas vacias
    # db = db.fillna("")
    #
    # # crear una nueva columna con compositor y autor juntos
    # db['compositor_autor'] = db.apply(
    #     lambda row: concaternar_autores(row['compositor'], row['autor']), axis=1)
    #
    # # Cambiar el formato de fecha de DD/MM/YYYY a YYYY-MM-DD o YYYY-MM o YYYY
    # db['fecha'] = db['fecha'].apply(convert_date_format)
    #
    # # crear una columna con solo el año
    # db['fecha_ano'] = db['fecha'].apply(extract_year)
    #
    # # Quitar todo el apellido en mayusculas
    # db['artista'] = db['artista'].apply(capitalize_uppercase_words)
    #
    # # CONVERTIR LOS NUMEROS A STRINGS Y QUITAR MAYUSCULAS Y ACENTOS
    # db['titulo_min'] = db['titulo'].apply(lambda x: convert_numbers_to_words(x) if pd.notna(x) else x)
    # db['artista_min'] = db['artista'].apply(lambda x: convert_numbers_to_words(x) if pd.notna(x) else x)
    # db['cantor_min'] = db['cantor'].apply(lambda x: convert_numbers_to_words(x) if pd.notna(x) else x)
    #
    # # palabras_mas_comunes(db,'artista_min')
    # db.to_csv(os.path.join(data_folder, "db.csv"),index=False,sep=';')
    return None


def concaternar_autores(compositor, autor):
    # if not isinstance(compositor, str):
    #     compositor = "?"
    # if not isinstance(autor, str):
    #     autor = "?"
    if compositor != "" and autor != "":
        concatenacion = "Musica: " + compositor + " - Letra: " + autor
    elif compositor != "":
        concatenacion = "Musica: " + compositor
    elif autor != "":
        concatenacion = "Letra: " + autor
    else:
        concatenacion = "Desconocido"

    return concatenacion


def adjust_text_size(text_frame, max_width_cm, max_font_size=100, min_font_size=10):
    # Convertir las medidas de cm a puntos (pt)
    max_width_pt = max_width_cm / 11500

    font_size = max_font_size
    text = ''.join(run.text for paragraph in text_frame.paragraphs for run in paragraph.runs)
    text_length = len(text)

    while font_size >= min_font_size:
        # Aplicar el tamaño de la fuente
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                run.font.name = DEFAULT_FONT_NAME

        # Cálculo del ancho estimado del texto basado en el tamaño de la fuente y la longitud del texto
        char_width = (DEFAULT_CHAR_WIDTH / 100) * font_size
        estimated_text_width = char_width * text_length

        # Si el ancho estimado del texto es menor o igual al ancho máximo permitido, se rompe el bucle
        if estimated_text_width <= max_width_pt:
            break

        font_size -= 1

    # Aplicar el tamaño de fuente final
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.name = DEFAULT_FONT_NAME