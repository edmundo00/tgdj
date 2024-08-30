import pandas as pd
from tinytag import TinyTag
import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter import ttk
from PIL import Image as PilImage, ImageTk
from pptx.util import Cm, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from unidecode import unidecode
import os
from os.path import join
from src.config.config import dropbox_path, image_folder, m3u_start_path, background_image_path, data_folder, output_folder, orchestra_folder, background_image_folder, merged_images_folder, DEFAULT_FONT_NAME, background_tango_degradado
from src.utils.utils import extract_year, separar_artistas, obtener_autores, convertir_segundos
from src.utils.funciones_para_diapos import *
# add_text_to_slide, calculate_positions, adjust_text_size
from datetime import datetime, timedelta
from src.utils.calcular_ancho_fuentes import FontWidthCalculator

# Estructura de self.result:
# --------------------------
# La variable `self.result` es un DataFrame de Pandas que almacena la estructura de la lista de reproducción
# agrupada por género y orquesta. Esta variable se construye en la función `create_structure()` y contiene
# las siguientes columnas:
#
# 1. `unique_value`: El valor único del género de cada grupo.
# 2. `repetition_count`: La cantidad de canciones que pertenecen al mismo grupo de género.
# 3. `position`: La posición del primer elemento del grupo en la lista original.
# 4. `same_orchestra`: Un valor booleano que indica si todas las canciones en el grupo tienen la misma orquesta.
# 5. `orchestra_value`: El nombre de la orquesta si todas las canciones en el grupo son de la misma orquesta;
#    de lo contrario, queda vacío.
# 6. `group_data`: Una lista de DataFrames, donde cada DataFrame contiene las filas de canciones que pertenecen
#    al mismo grupo, manteniendo la estructura original de la lista de reproducción.

# Ejemplo de self.result:
# ------------------------
# unique_value   repetition_count   position   same_orchestra   orchestra_value   group_data
# tango          5                  0          True             Orquesta XYZ      DataFrame (5 filas)
# tango vals     3                  5          True             Orquesta ABC      DataFrame (3 filas)
# tango milonga  4                  8          False                              DataFrame (4 filas)

# Ejemplo de groupdata

# [1 rows x 13 columns],                   title         artist1  ...   ano       artist2
# 0              El olivo  Juan D'Arienzo  ...  1941  Héctor Mauré
# 1  Si la llegaran a ver  Juan D'Arienzo  ...  1943  Héctor Mauré
# 2           Humillación  Juan D'Arienzo  ...  1941  Héctor Mauré
# 3               Amarras  Juan D'Arienzo  ...  1944  Héctor Mauré

# Las columnas estan definidas en read_audio_tags
# columns = ['title', 'artist1', 'artist2', 'album', 'year', 'genre', 'composer', 'lyrics', 'bpm', 'duration',
#                        'extension', 'bitrate']

# Descripción de la función canciones_tanda:
# -------------------------------------------
# La función `canciones_tanda(tanda, tags)` toma como argumentos el número de tanda (`tanda`) y una lista
# de etiquetas (`tags`) que especifican los metadatos que se desean extraer. La función devuelve una lista
# de listas, donde cada sublista contiene los valores de las etiquetas especificadas para una canción en
# la tanda seleccionada.
#
# Funcionamiento:
# - La función verifica si la tanda solicitada existe en `self.result`.
# - Si la tanda existe, extrae el DataFrame correspondiente desde la columna `group_data` de `self.result`.
# - Luego, recorre cada fila del DataFrame y extrae los valores de las etiquetas especificadas.
# - Devuelve una lista de listas, donde cada sublista contiene los valores de las etiquetas para una canción.

# Ejemplo de uso de canciones_tanda:
# -----------------------------------
# Si se llama `canciones_tanda(1, ['title', 'ano', 'composer'])` para obtener los títulos, años y compositores
# de las canciones en la primera tanda:
#
# Entrada:
# tanda = 1
# tags = ['title', 'ano', 'composer']
#
# Salida:
# [
#   ['La Cumparsita', '1916', 'Gerardo Matos Rodríguez'],
#   ['El Choclo', '1903', 'Ángel Villoldo'],
#   ['A Media Luz', '1925', 'Edgardo Donato'],
#   ...
# ]
#
# Cada sublista contiene los valores de 'title', 'ano' y 'composer' para una canción en la tanda especificada.



class PresentationApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Presentation Creator")
        self.root.geometry('1500x800')  # Define el tamaño inicial de la ventana

        # Uso de la clase FontWidthCalculator
        self.calculadora = FontWidthCalculator()
        # Crear la base de datos y cargarla en memoria
        self.calculadora.crear_base_de_datos()

        # self.root.state('zoomed')
        # Set the window to be on top
        # self.root.attributes('-topmost', True)
        # self.root.after_idle(lambda: self.root.attributes('-topmost', False))  # Allow interaction with other windows

        # Default paths

        output_folder = join(dropbox_path, "MUSICA", "MP3", "TANGO", "other_stuff", "Presentacion")


        self.merged_image_folder = image_folder
        self.m3u_file_path = None
        self.audio_files = []
        # self.audio_tags = []

        # Create Menu
        menubar = tk.Menu(root)
        root.config(menu=menubar)

        file_menu = tk.Menu(menubar, tearoff=0)
        file_menu.add_command(label="Open", command=self.open_m3u_file)
        menubar.add_cascade(label="File", menu=file_menu)

        preferences_menu = tk.Menu(menubar, tearoff=0)
        preferences_menu.add_command(label="Preferences", command=self.open_preferences)
        menubar.add_cascade(label="Preferences", menu=preferences_menu)



        # Entrada del Nombre de la Milonga
        tk.Label(root, text="Nombre de la Milonga:").grid(row=0, column=0, padx=10, pady=10, sticky="w")
        self.nombre_milonga_entry = tk.Entry(root, width=50)
        self.nombre_milonga_entry.grid(row=0, column=1, columnspan=3, padx=10, pady=10, sticky="ew")
        self.nombre_milonga_entry.insert(0, "Milonga de la Fuente")

        # Entrada de la Fecha
        tk.Label(root, text="Fecha:").grid(row=1, column=0, padx=10, pady=10, sticky="w")
        self.fecha_entry = tk.Entry(root, width=50)
        self.fecha_entry.grid(row=1, column=1, columnspan=3, padx=10, pady=10, sticky="ew")
        self.fecha_entry.insert(0, "24 de Agosto de 2024")

        # Entradas de Hora de Inicio, Finalización y Duración en la misma fila
        tk.Label(root, text="Inicio:").grid(row=2, column=0, padx=10, pady=10, sticky="w")
        self.hora_inicio_entry = tk.Entry(root, width=10)
        self.hora_inicio_entry.grid(row=2, column=1, padx=10, pady=10)
        self.hora_inicio_entry.insert(0, "21:00")  # Hora por defecto

        tk.Label(root, text="Final:").grid(row=2, column=2, padx=10, pady=10, sticky="w")
        self.hora_fin_entry = tk.Entry(root, width=10)
        self.hora_fin_entry.grid(row=2, column=3, padx=10, pady=10)
        self.hora_fin_entry.insert(0, "00:00")  # Hora por defecto

        tk.Label(root, text="Duración:").grid(row=2, column=4, padx=10, pady=10, sticky="w")
        self.duracion_label = tk.Label(root, text="", width=20, bg="white", relief="sunken", anchor='w', padx=2, pady=2)
        self.duracion_label.grid(row=2, column=5, padx=10, pady=10)

        # Vincular eventos de cambio para actualizar la duración automáticamente
        self.hora_inicio_entry.bind("<KeyRelease>", self.calcular_duracion)
        self.hora_fin_entry.bind("<KeyRelease>", self.calcular_duracion)

        # Mostrar la duración inicial basada en los valores por defecto
        self.calcular_duracion()


        # Thumbnail de fondo
        self.update_background_thumbnail()

        # Botón para cargar el archivo M3U
        load_m3u_button = ttk.Button(root, text="Load M3U File", command=self.open_m3u_file)
        load_m3u_button.grid(row=3, column=0, columnspan=6, pady=10, sticky="ew")

        # Botón para crear la presentación
        self.create_button = ttk.Button(root, text="Crear Presentación", command=self.create_presentation,
                                        state=tk.DISABLED)
        self.create_button.grid(row=5, column=0, columnspan=6, pady=20, sticky="nsew")

        # Configuración del árbol de vista de títulos, artistas y géneros
        self.tree = ttk.Treeview(root, columns=("Title", "Artist", "Genre"), show='headings')
        self.tree.heading("Title", text="Title")
        self.tree.heading("Artist", text="Artist")
        self.tree.heading("Genre", text="Genre")
        self.tree.grid(row=4, column=0, columnspan=6, padx=10, pady=10, sticky="nsew")

        # Configuración del frame de la estructura
        self.estructura_frame = ttk.LabelFrame(root, text="Estructura")
        self.estructura_frame.grid(row=0, column=6, rowspan=6, padx=10, pady=10, sticky="nsew")

        self.estructura_tree = ttk.Treeview(self.estructura_frame, columns=("Col1", "Col2", "Col3", "Col4"),
                                            show='headings')
        self.estructura_tree.heading("Col1", text="Genero")
        self.estructura_tree.heading("Col2", text="Orquesta")
        self.estructura_tree.heading("Col3", text="Columna 3")
        self.estructura_tree.heading("Col4", text="Numero de canciones")
        self.estructura_tree.pack(expand=True, fill='both', padx=10, pady=10)

        # Configuración de las filas y columnas
        root.grid_rowconfigure(4, weight=1)
        root.grid_columnconfigure(1, weight=1)
        root.grid_columnconfigure(2, weight=1)

    def calcular_duracion(self, event=None):
        formato_hora = "%H:%M"
        try:
            hora_inicio = datetime.strptime(self.hora_inicio_entry.get(), formato_hora)
            hora_fin = datetime.strptime(self.hora_fin_entry.get(), formato_hora)

            # Si la hora de fin es menor o igual que la de inicio, se considera al día siguiente
            if hora_fin <= hora_inicio:
                hora_fin += timedelta(days=1)

            duracion = hora_fin - hora_inicio
            horas, minutos = divmod(duracion.seconds, 3600)
            minutos //= 60

            # Mostrar la duración en la etiqueta
            self.duracion_label.config(text=f"{horas} horas y {minutos} minutos")

        except ValueError:
            self.duracion_label.config(text="Formato inválido")




    def create_structure(self):
        # Check if 'genre' column exists and is not empty
        if 'genre' not in self.df.columns or self.df['genre'].empty:
            print("The 'genre' column is missing or empty.")
            return

        # Original list
        original_list = self.df['genre']

        # Create a new DataFrame from the 'genre' and 'artist1' columns
        df_genres = pd.DataFrame({'genre': original_list, 'artist1': self.df['artist1']})

        # Identify when the genre value changes
        df_genres['change'] = df_genres['genre'].ne(df_genres['genre'].shift())

        # Create a group identifier based on consecutive identical genre values
        df_genres['group'] = df_genres['change'].cumsum()

        # Check if all orchestras in the same group are the same
        df_genres['same_orchestra'] = df_genres.groupby('group')['artist1'].transform(lambda x: x.nunique() == 1)

        # Collect DataFrames for each group in a list
        group_data_list = [self.df.loc[indices].reset_index(drop=True) for group, indices in
                           df_genres.groupby('group').groups.items()]

        # Group by the group identifier and calculate the unique values, counts, positions, same_orchestra status, and store the DataFrame for each group
        self.result = df_genres.groupby('group').agg(
            unique_value=('genre', 'first'),
            repetition_count=('genre', 'size'),
            position=('group', 'idxmin'),
            same_orchestra=('same_orchestra', 'first'),
            orchestra_value=('artist1', lambda x: x.iloc[0] if x.nunique() == 1 else "")
        ).reset_index(drop=True)


        # Attach the group DataFrames to the result DataFrame
        self.result['group_data'] = group_data_list

        # Lista de excepciones donde no se debe añadir "Orquesta de"
        excepciones = ["Orquesta", "Quinteto", "Sexteto", "Hugo Diaz"]

        # Añadir la columna 'titulo_orquesta' según las condiciones dadas
        self.result['titulo_orquesta'] = self.result['orchestra_value'].apply(
            lambda orchestra_value: orchestra_value
            if any(orchestra_value.startswith(excepcion) for excepcion in excepciones)
            else f"Orquesta de {orchestra_value}"
        )
        self.result['genero_autores'] = [
            self.obtener_genero_autores(
                row['unique_value'],
                obtener_autores(self.canciones_tanda(tanda_number + 1, ['artist2']))
            )
            for tanda_number, row in self.result.iterrows()
        ]

        # Añadir la columna 'total_duration' con la suma de 'duration' en cada group_data
        self.result['duracion_total'] = self.result['group_data'].apply(
            lambda group_df: group_df['duration'].sum() if 'duration' in group_df.columns else 0
        )

        # Clear the current content of the estructura_tree
        for i in self.estructura_tree.get_children():
            self.estructura_tree.delete(i)

        # Configure tag styles for the Treeview
        self.estructura_tree.tag_configure("SameOrchestra.Treeview", background="honeydew", foreground="black")
        self.estructura_tree.tag_configure("DifferentOrchestra.Treeview", background="red", foreground="white")

        # Insert new rows into the estructura_tree with the new column data
        for _, row in self.result.iterrows():

            if row['same_orchestra']:
                row_style = "SameOrchestra.Treeview"
            else:
                row_style = "DifferentOrchestra.Treeview"

            # Here you could do something with row['group_data'] if needed

            self.estructura_tree.insert("", "end", values=(
                row['unique_value'], row['orchestra_value'], '', row['repetition_count'], row['same_orchestra']),
                                        tags=(row_style,))




        #
        # # Save the result DataFrame to an CSV file
        # save_csv_path = os.path.join(data_folder, "self_result.csv")
        # self.result.to_csv(save_csv_path, index=False)


    def canciones_tanda(self, tanda, tags):
        # Check if the result has at least 'tanda' rows
        if len(self.result) < tanda:
            print(f"The result DataFrame has less than {tanda} rows.")
            return None

        # Access the specified row (index tanda-1 since indexing is 0-based)
        selected_row = self.result.iloc[tanda - 1]

        # Extract the DataFrame from the group_data column
        group_df = selected_row['group_data']

        # Initialize a list to store the tags for each row
        list_of_values = []

        # Loop through each row in group_df
        for index, row in group_df.iterrows():
            # Extract the values of the specified tags in the current row
            selected_values = [row[tag] for tag in tags]
            # Append the list of values to the list_of_values
            list_of_values.append(selected_values)

        return list_of_values

    def open_m3u_file(self):
        # Temporarily make the window non-topmost to allow file dialog interaction
        # self.root.attributes('-topmost', False)

        # Open file dialog
        self.m3u_file_path = filedialog.askopenfilename(
            initialdir=m3u_start_path,
            filetypes=[("M3U files", "*.m3u"), ("All files", "*.*")]
        )

        # # Return focus to the popup window
        # self.root.attributes('-topmost', True)

        if self.m3u_file_path:
            # Clear previous data
            if hasattr(self, 'df'):
                self.df = self.df.iloc[0:0]  # Clear the DataFrame but keep the columns

            self.audio_files = []  # Clear the audio files list
            # self.audio_tags = []  # Clear the audio tags list

            # Load the M3U file and extract data
            with open(self.m3u_file_path, 'r', encoding='utf-8') as file:
                for line in file:
                    line = line.strip()
                    if line and not line.startswith("#"):
                        if os.path.exists(line):
                            self.audio_files.append(line)
                            tags = self.read_audio_tags(line)
                            # self.audio_tags.append(tags)
                        else:
                            modified_path = dropbox_path + line.split("Dropbox", 1)[1]
                            if os.path.exists(modified_path):
                                self.audio_files.append(modified_path)
                                tags = self.read_audio_tags(modified_path)
                                # self.audio_tags.append(tags)

            self.create_structure()

            if self.audio_files:
                self.display_m3u_summary()
                self.create_button.config(state=tk.NORMAL)

    def read_audio_tags(self, file_path):
        # Check if the DataFrame is already initialized, if not, initialize it
        if not hasattr(self, 'df'):
            columns = ['title', 'artist1', 'artist2', 'album', 'year', 'genre', 'composer', 'lyrics', 'bpm', 'duration',
                       'extension', 'bitrate']
            self.df = pd.DataFrame(columns=columns)

        try:
            # Extract tags from the audio file
            tags = TinyTag.get(file_path)
            self.artists1, self.artists2 = separar_artistas(tags.artist)

            # Extract file extension and bitrate
            extension = file_path.split('.')[-1]
            bitrate = tags.bitrate if tags.bitrate else "Unknown Bitrate"

            # Prepare a dictionary with the extracted tags
            tag_data = {
                "title": tags.title if tags.title else "Unknown Title",
                "artist1": self.artists1 if tags.artist else "Unknown Artist",
                "artist2": self.artists2 if tags.artist else "Unknown Artist",
                "album": tags.album if tags.album else "Unknown Album",
                "year": tags.year if tags.year else "Unknown Year",
                "genre": tags.genre if tags.genre else "Unknown Genre",
                "composer": tags.composer if tags.composer else "Unknown Composer",
                "lyrics": tags.lyrics if hasattr(tags, 'lyrics') else "Unknown Lyrics",
                "bpm": tags.bpm if hasattr(tags, 'bpm') else "Unknown BPM",
                "duration": tags.duration if tags.duration else "Unknown Duration",
                "extension": extension,
                "bitrate": bitrate
            }

            # Remove entries in tag_data where the value is None, NaN, or an empty string
            tag_data = {key: (value if value else "Unknown") for key, value in tag_data.items() if
                        value not in [None, ""]}

        except Exception as e:
            # Fallback tags in case of an error
            tag_data = {
                "title": "Unknown Title",
                "artist1": "Unknown Artist",
                "artist2": "Unknown Artist",
                "album": "Unknown Album",
                "year": "Unknown Year",
                "genre": "Unknown Genre",
                "composer": "Unknown Composer",
                "lyrics": "Unknown Lyrics",
                "bpm": "Unknown BPM",
                "duration": "Unknown Duration",
                "extension": extension if 'extension' in locals() else "Unknown Extension",
                "bitrate": bitrate if 'bitrate' in locals() else "Unknown Bitrate"
            }

        # Convert tag_data to a DataFrame and replace empty or all-NA columns
        tag_df = pd.DataFrame([tag_data])

        # Remove any columns that are all NA or have only "Unknown" values before concatenating
        tag_df.replace("Unknown", pd.NA, inplace=True)
        tag_df.dropna(axis=1, how='all', inplace=True)

        # If self.df is empty, assign tag_df directly, otherwise concatenate
        if self.df.empty:
            self.df = tag_df
        else:
            self.df = pd.concat([self.df, tag_df], ignore_index=True)

        # Apply the extract_year function to the 'year' column and create a new 'ano' column
        self.df['ano'] = self.df['year'].apply(extract_year)

    def display_m3u_summary(self):
        for i in self.tree.get_children():
            self.tree.delete(i)

        for index, row in self.df.iterrows():
            self.tree.insert("", "end", values=(row.title, row.artist1, row.genre))

    def open_preferences(self):
        pref_window = tk.Toplevel(self.root)
        pref_window.title("Preferences")
        pref_window.geometry("600x200")

        tk.Label(pref_window, text="Output Path:").grid(row=0, column=0, padx=10, pady=10, sticky="w")
        output_path_entry = tk.Entry(pref_window, width=70)
        output_path_entry.grid(row=0, column=1, padx=10, pady=10)
        output_path_entry.insert(0, output_folder)

        tk.Label(pref_window, text="Background Image:").grid(row=1, column=0, padx=10, pady=10, sticky="w")
        background_entry = tk.Entry(pref_window, width=70)
        background_entry.grid(row=1, column=1, padx=10, pady=10)
        background_entry.insert(0, background_image_path)

        tk.Label(pref_window, text="M3U Start Path:").grid(row=2, column=0, padx=10, pady=10, sticky="w")
        m3u_start_path_entry = tk.Entry(pref_window, width=70)
        m3u_start_path_entry.grid(row=2, column=1, padx=10, pady=10)
        m3u_start_path_entry.insert(0, m3u_start_path)

        def save_preferences():
            output_folder = output_path_entry.get()
            background_image_path = background_entry.get()
            m3u_start_path = m3u_start_path_entry.get()
            messagebox.showinfo("Preferences", "Preferences saved!")
            pref_window.destroy()

        save_button = ttk.Button(pref_window, text="Save", command=save_preferences)
        save_button.grid(row=3, column=0, columnspan=2, pady=20)

    def update_background_thumbnail(self):
        image = PilImage.open(background_image_path)
        image.thumbnail((100, 100))
        self.thumbnail_image = ImageTk.PhotoImage(image)

        if hasattr(self, 'thumbnail_button'):
            self.thumbnail_button.config(image=self.thumbnail_image)
        else:
            self.thumbnail_button = tk.Button(self.root, image=self.thumbnail_image,
                                              command=self.select_background_image)
            self.thumbnail_button.grid(row=0, column=4, columnspan=2, pady=10)

    def select_background_image(self):
        new_image_path = filedialog.askopenfilename(
            initialdir=os.path.dirname(background_image_path),
            filetypes=[("Image files", "*.jpg;*.jpeg;*.png;*.webp;*.bmp"), ("All files", "*.*")]
        )
        if new_image_path:
            background_image_path = new_image_path
            self.update_background_thumbnail()

    def apply_gradient_overlay(self):
        # Open the original image
        base = PilImage.open(background_image_path).convert("RGBA")

        # #orchestra = PilImage.open(self.orchestra_path).convert("RGBA")
        #
        # maximum_in_width = 0.4 * base.width
        # maximum_in_height = base.height
        #
        # resize_factor = maximum_in_height / orchestra.height
        #
        # if (resize_factor * orchestra.width) > maximum_in_width:
        #     resize_factor = maximum_in_width / orchestra.width
        #
        # new_width = int(orchestra.width * resize_factor)
        # new_height = int(orchestra.height * resize_factor)
        #
        # # Resize the orchestra image
        # orchestra = orchestra.resize((new_width, new_height), PilImage.Resampling.LANCZOS)

        # Create a gradient overlay
        gradient = PilImage.new('L', (base.width, base.height))
        for x in range(base.width):
            # Gradient from opaque dark gray (on the left) to fully transparent (on the right)
            gradient_value = int(
                255 * (1 - x / base.width))  # Fully opaque on the left to fully transparent on the right
            for y in range(base.height):
                gradient.putpixel((x, y), gradient_value)

        # Convert gradient to RGBA
        gradient_rgba = PilImage.new('RGBA', base.size)
        for y in range(base.height):
            for x in range(base.width):
                gradient_rgba.putpixel((x, y), (
                105, 105, 105, gradient.getpixel((x, y))))  # Dark gray with varying transparency

        # Apply the gradient to the base image
        combined = PilImage.alpha_composite(base, gradient_rgba)

        # Paste the resized orchestra image onto the combined image
        # position = (0, combined.height - orchestra.height)  # Centering
        # combined.paste(orchestra, position, orchestra)  # The third argument is the mask to maintain transparency

        # Save the gradient and the combined image
        # gradient_rgba.save(self.path_gradiente, "PNG")
        # self.merged_image_path = join(image_folder, "merged_background.png")

        self.merged_image_path = join(background_image_folder, "background_tango_degradado.png")

        combined.save(self.merged_image_path, "PNG")

    def create_slide_for_tanda(self, prs, tanda_number, titulo, titulo_orquesta, subtitulo, genero, lista_canciones, positions_initial):

        # self.apply_gradient_overlay()

         # NAME OF THE ORCHESTRA
        orchestra_value = titulo
        orchestra_value_min = unidecode(orchestra_value).lower()
        self.tanda_gender = genero
        lista_canciones = lista_canciones

        # Determine the background image based on genre
        if 'tango' not in self.tanda_gender:
            self.merged_image_path = join(background_image_folder, "background_cortina.png")
        else:
            # self.merged_image_path = join(merged_images_folder, f'{orchestra_value_min}_background.png')
            self.merged_image_path = join(background_image_folder, "background_tango_degradado.png")

        # Add a slide with a title and content layout
        slide_layout = prs.slide_layouts[5]  # Use a blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Set the merged background image with gradient
        slide.shapes.add_picture(self.merged_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)


        img_path = join(image_folder, 'orquestas', f'{orchestra_value_min}.png')

        if os.path.exists(img_path):
            add_resized_image_to_slide(slide,img_path, positions_initial["maxima_anchura_image"], prs)

        positions_calculated = transformar_posiciones(prs, positions_initial)


        if 'tango' not in self.tanda_gender:
            title, year, composer = lista_canciones[0]

            add_text_to_slide(
                slide,
                self.calculadora,
                f'{orchestra_value}',
                positions_calculated["cortina_title"],
                positions_calculated["offset_shadow"]*2.5,
                positions_initial["fuentes"]["cortina_titulo"]['tamaño'],
                positions_initial["fuentes"]["cortina_titulo"]['tipo_fuente'],
                (255, 255, 255),  # Color blanco
                True,  # Negrita activada
                True,  # Sombra activada
                adjust_size = False
            )


            add_text_to_slide(
                slide,
                self.calculadora,
                f'{title}   ({year})',
                positions_calculated["cortina_subtitle"],
                positions_calculated["offset_shadow"]*2.5,
                positions_initial["fuentes"]["cortina_artista"]['tamaño'],
                positions_initial["fuentes"]["cortina_artista"]['tipo_fuente'],
                (255, 255, 255),  # Color blanco
                True,  # Negrita activada
                True,  # Sombra activada
                adjust_size = False,
            )

        else:

            add_text_to_slide(
                slide,
                self.calculadora,
                titulo_orquesta,
                positions_calculated["tanda_orquesta_shadow"],
                positions_calculated["offset_shadow"]*2.5,
                positions_initial["fuentes"]["orquesta"]['tamaño'],
                positions_initial["fuentes"]["orquesta"]['tipo_fuente'],
                (255, 255, 255),  # Color blanco
                True,  # Negrita activada
                True,  # Sombra activada
                adjust_size = False
            )

            add_text_to_slide(
                slide,
                self.calculadora,
                subtitulo,
                positions_calculated["tanda_cantor_shadow"],
                positions_calculated["offset_shadow"],
                positions_initial["fuentes"]["estilo"]['tamaño'],
                positions_initial["fuentes"]["estilo"]['tipo_fuente'],
                (255, 255, 255),  # Color blanco
                True,  # Negrita activada
                True,  # Sombra activada
                adjust_size=False
            )

            add_text_to_slide(
                slide,
                self.calculadora,
                f'© TDJ Edmundo Fraga\n{self.nombre_milonga_entry.get()}\n{self.fecha_entry.get()}',
                positions_calculated["firma_tgdj_box"],
                positions_calculated["offset_shadow"],
                positions_initial["fuentes"]["firma"]['tamaño'],
                positions_initial["fuentes"]["firma"]['tipo_fuente'],
                (255, 255, 255),  # Color blanco
                False,  # Negrita activada
                True,  # Sombra activada
                adjust_size=False,
                border_color_rgb = (0, 0, 0),
                border_width_pt = 2
            )

            linea_divisoria = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                                         positions_calculated["linea_divisoria"][0],
                                                         positions_calculated["linea_divisoria"][1],
                                                         positions_calculated["linea_divisoria"][0] +
                                                         positions_calculated["linea_divisoria"][2],
                                                         positions_calculated["linea_divisoria"][1] +
                                                         positions_calculated["linea_divisoria"][3])
            linea_divisoria.line.color.rgb = RGBColor(255, 255, 255)
            linea_divisoria.line.width = Pt(7)

            for rows in lista_canciones:
                title, year, composer = rows

                add_text_to_slide(
                    slide,
                    self.calculadora,
                    title,
                    positions_calculated["canciones_start"],
                    positions_calculated["offset_shadow"],
                    positions_initial["fuentes"]["canciones"]['tamaño'],
                    positions_initial["fuentes"]["canciones"]['tipo_fuente'],
                    (255, 255, 255),  # Color blanco
                    False,  # Negrita activada
                    True,  # Sombra activada
                    adjust_size=False,
                    extra_paragraph_text = composer,
                    extra_run_text = f'   ({year})',
                    extra_paragraph_settings ={'font_name':positions_initial["fuentes"]["autores"]['tipo_fuente'], 'tamano_fuente':positions_initial["fuentes"]["autores"]['tamaño'], 'font_color_rgb':RGBColor(255, 255, 255), 'is_bold':False, 'is_italic':True},
                    extra_run_settings = {'font_name':positions_initial["fuentes"]["anos"]['tipo_fuente'], 'tamano_fuente':positions_initial["fuentes"]["anos"]['tamaño'], 'font_color_rgb':RGBColor(255, 255, 255), 'is_bold':False, 'is_italic':False},
                )

                positions_calculated["canciones_start"][1] = positions_calculated["canciones_start"][1] + positions_calculated["canciones_spacing"]

    def create_first_slide(self, prs, posiciones, datos):
        # "Milonga": {"left": Cm(1.6), "top": Cm(0), "right": Cm(1), "height": Cm(3)},
        # "Fecha_Milonga": {"left": Cm(3), "top": Cm(2.3), "right": Cm(1), "height": Cm(2)},
        # "Duracion_Milonga": {"left": Cm(3), "top": Cm(2.3), "right": Cm(1), "height": Cm(2)},
        # "firma": {"left": Cm(1), "top": Cm(15.5), "right": Cm(24), "height": Cm(2.8)},
        # "linea_divisoria": {"left": Cm(13.5), "top": Cm(4.5), "right": Cm(1), "height": Cm(0)},
        # "tandas_start": {"left": Cm(4), "top": Cm(4), "right": Cm(1), "height": Cm(15)},
        # "tandas_spacing": Cm(0.8),
        # "offset_shadow": Cm(0.06),
        # "maxima_anchura_image": 1,
        # "fuentes": {
        #     "Milonga": {"tamaño": 100, "tipo_fuente": "Arial"},
        #     "Fecha_Milonga": {"tamaño": 50, "tipo_fuente": "Arial"},
        #     "Duracion_Milonga": {"tamaño": 55, "tipo_fuente": "Arial"},
        #     "tandas": {"tamaño": 20, "tipo_fuente": "Bernard MT Condensed"},
        #     "duracion_estimada": {"tamaño": 20, "tipo_fuente": "Arial"},
        #     "firma": {"tamaño": 20, "tipo_fuente": "Arial"}
        # datos = {'dj': 'Edmundo Fraga',
        #          'milonga': self.nombre_milonga_entry.get(),
        #          'fecha': self.fecha_entry.get(),
        #          'inicio':self.hora_inicio_entry.get(),
        #          'finalizacion':self.hora_fin_entry.get(),
        #          'duracion': self.duracion_label.cget("text"),
        #          'tandas': lista_texto,
        #          'duracion_total_sin_cortina' : tx_duracion_total_sin_cortina,
        #          'duracion_total_estimada': tx_duracion_total_estimado
        #          }

        # Add a slide with a title and content layout
        slide_layout = prs.slide_layouts[5]  # Use a blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Set the merged background image with gradient
        slide.shapes.add_picture(background_tango_degradado, 0, 0, width=prs.slide_width, height=prs.slide_height)

        pos = transformar_posiciones(prs, posiciones)

        add_text_to_slide(
            slide,
            self.calculadora,
            datos["milonga"],
            pos["Milonga"],
            pos["offset_shadow"],
            pos['fuentes']['Milonga']['tamaño'],
            pos['fuentes']['Milonga']['tipo_fuente'],
            (255, 255, 255),  # Color blanco
            False,  # Negrita activada
            True,  # Sombra activada
        )

        add_text_to_slide(
            slide,
            self.calculadora,
            datos["fecha"] + " de " + datos["inicio"] +" a "+  datos["finalizacion"] + " duracion: " + datos["duracion"],
            pos["Fecha_Milonga"],
            pos["offset_shadow"],
            pos['fuentes']['Fecha_Milonga']['tamaño'],
            pos['fuentes']['Fecha_Milonga']['tipo_fuente'],
            (255, 255, 255),  # Color blanco
            False,  # Negrita activada
            True,  # Sombra activada
            adjust_size=False,
        )

        add_text_to_slide(
            slide,
            self.calculadora,
            datos['duracion_total_sin_cortina'] + ", " + datos['duracion_total_estimada'],
            pos["Duracion_Milonga"],
            pos["offset_shadow"],
            pos['fuentes']['Duracion_Milonga']['tamaño'],
            pos['fuentes']['Duracion_Milonga']['tipo_fuente'],
            (255, 255, 255),  # Color blanco
            False,  # Negrita activada
            True,  # Sombra activada
            adjust_size=False,
        )

        for row in datos['tandas']:
            add_text_to_slide(
                slide,
                self.calculadora,
                row,
                pos["tandas_start"],
                pos["offset_shadow"],
                pos['fuentes']['tandas']['tamaño'],
                pos['fuentes']['tandas']['tipo_fuente'],
                (255, 255, 255),  # Color blanco
                False,  # Negrita activada
                True,  # Sombra activada
                adjust_size=False,
            )

            pos["tandas_start"][1] = pos["tandas_start"][1] + pos["tandas_spacing"]

    def preparar_datos_first_slide(self):

        tandas_sin_cortinas = self.result[self.result['unique_value'].str.contains('tango', case=False, na=False)].copy()
        tc=tandas_sin_cortinas
        lista_texto = []
        duracion_total_sin_cortina = 0
        counter = 1
        for index, row in tc.iterrows():
            duracion = convertir_segundos(row['duracion_total'], "x minutos y x segundos")
            columna =(f"{counter} - {row['titulo_orquesta']} {row['genero_autores']}, {row['repetition_count']} temas con una duración de {duracion}")
            lista_texto.append(columna)
            duracion_total_sin_cortina += row['duracion_total']
            counter+=1

        duracion_cortinas_extimado = (len(tc)-1)*60
        duracion_total_estimado = duracion_cortinas_extimado + duracion_total_sin_cortina
        tx_duracion_total_sin_cortina = f'Duración total sin cortinas: {convertir_segundos(duracion_total_sin_cortina, "x minutos y x segundos")}'
        tx_duracion_cortinas_extimado = convertir_segundos(duracion_cortinas_extimado, "x minutos y x segundos")
        tx_duracion_total_estimado = f'Duración total estimada: {convertir_segundos(duracion_total_estimado, "x minutos y x segundos")}'

        datos = {'dj': 'Edmundo Fraga',
                 'milonga': self.nombre_milonga_entry.get(),
                 'fecha': self.fecha_entry.get(),
                 'inicio':self.hora_inicio_entry.get(),
                 'finalizacion':self.hora_fin_entry.get(),
                 'duracion': self.duracion_label.cget("text"),
                 'tandas': lista_texto,
                 'duracion_total_sin_cortina' : tx_duracion_total_sin_cortina,
                 'duracion_total_estimada': tx_duracion_total_estimado
                 }

        return  datos


    def obtener_genero_autores(self, genero, autores):
        if 'vals' in genero:
            subtitulo = "Vals "
            if autores == 'instrumental':
                subtitulo = subtitulo + 'instrumental'
            else:
                subtitulo = subtitulo + 'con ' + autores
        elif 'milonga' in genero:
            subtitulo = "Milonga "
            if autores == 'instrumental':
                subtitulo = subtitulo + 'instrumental'
            else:
                subtitulo = subtitulo + 'con ' + autores
        elif genero == 'tango':
            if autores == 'instrumental':
                subtitulo = 'Instrumental'
            else:
                subtitulo = 'Con ' + autores
        else:
            subtitulo = ''
        return subtitulo


    def create_presentation(self):

        nombre_milonga = self.nombre_milonga_entry.get()
        fecha = self.fecha_entry.get()

        if not nombre_milonga or not fecha:
            messagebox.showwarning("Input Error", "Please fill in both 'Nombre de la Milonga' and 'Fecha'.")
            return

        if not self.m3u_file_path or not self.audio_files:
            messagebox.showwarning("Input Error", "Please load an M3U file.")
            return

        # Path to save the presentation
        output_file = join(output_folder, "presentation.pptx")

        # Create a new PowerPoint presentation
        prs = Presentation()

        # Set slide dimensions (7:4 aspect ratio)
        prs.slide_width = Cm(33.87)
        prs.slide_height = Cm(19.05)

        posiciones_first_slide = {
                "Milonga": {"left": Cm(1), "top": Cm(0), "right": Cm(1), "height": Cm(3)},
                "Fecha_Milonga": {"left": Cm(1.75), "top": Cm(2.5), "right": Cm(1), "height": Cm(2)},
                "Duracion_Milonga": {"left": Cm(1), "top": Cm(17.25), "right": Cm(1), "height": Cm(2)},
                "firma": {"left": Cm(1), "top": Cm(15.5), "right": Cm(24), "height": Cm(2.8)},
                "linea_divisoria": {"left": Cm(13.5), "top": Cm(4.5), "right": Cm(1), "height": Cm(0)},
                "tandas_start": {"left": Cm(3), "top": Cm(4), "right": Cm(1), "height": Cm(15)},
                "tandas_spacing" : Cm(0.7),
                "offset_shadow": Cm(0.06),
                "maxima_anchura_image": 1,
                "fuentes": {
                    "Milonga": {"tamaño": 60, "tipo_fuente": "Arial"},
                    "Fecha_Milonga": {"tamaño": 20, "tipo_fuente": "Arial"},
                    "Duracion_Milonga": {"tamaño": 16, "tipo_fuente": "Arial"},
                    "tandas": {"tamaño": 15, "tipo_fuente": "Bernard MT Condensed"},
                    "duracion_estimada": {"tamaño": 16, "tipo_fuente": "Arial"},
                    "firma": {"tamaño": 20, "tipo_fuente": "Arial"}
                }
        }

        datos_first_slide = self.preparar_datos_first_slide()

        self.create_first_slide(prs, posiciones_first_slide, datos_first_slide)

        # Add slides for different tanda numbers
        for tanda_number in range(1, self.result.shape[0]+1):  # Adjust the range as needed
        # for tanda_number in range(1, 5):  # Adjust the range as needed
            titulo = self.result.iloc[tanda_number - 1]['orchestra_value']
            titulo_orquesta = self.result.iloc[tanda_number - 1]['titulo_orquesta']
            genero = self.tanda_gender = self.result.iloc[tanda_number - 1]['unique_value']

            subtitulo = self.result.iloc[tanda_number - 1]['genero_autores']

            canciones = self.canciones_tanda(tanda_number, ['title', 'ano', 'composer'])

            positions_initial = {
                "cortina_title": {"left": Cm(5), "top": Cm(5), "right": Cm(1), "height": Cm(5)},
                "cortina_subtitle": {"left": Cm(8), "top": Cm(10), "right": Cm(1), "height": Cm(5)},
                "tanda_orquesta_shadow": {"left": Cm(1.6), "top": Cm(0), "right": Cm(1), "height": Cm(3)},
                "tanda_cantor_shadow": {"left": Cm(3), "top": Cm(2.3), "right": Cm(1), "height": Cm(2)},
                "firma_tgdj_box": {"left": Cm(1), "top": Cm(15.5), "right": Cm(24), "height": Cm(2.8)},
                "linea_divisoria": {"left": Cm(13.5), "top": Cm(4.5), "right": Cm(1), "height": Cm(0)},
                "canciones_start": {"left": Cm(13.5), "top": Cm(4.5), "right": Cm(1), "height": Cm(3)},
                "canciones_spacing": Cm(3.5),
                "offset_shadow": Cm(0.05),
                "maxima_anchura_image": 1,
                "fuentes": {
                "cortina_titulo": {"tamaño": 100, "tipo_fuente": DEFAULT_FONT_NAME},
                "cortina_artista": {"tamaño": 50, "tipo_fuente": DEFAULT_FONT_NAME},
                "orquesta": {"tamaño": 55, "tipo_fuente": DEFAULT_FONT_NAME},
                "estilo": {"tamaño": 35, "tipo_fuente": DEFAULT_FONT_NAME},
                "canciones": {"tamaño": 50, "tipo_fuente": 'Bernard MT Condensed'},
                "anos": {"tamaño": 20, "tipo_fuente": "Arial"},
                "autores": {"tamaño": 20, "tipo_fuente": 'Arial Narrow'},
                "firma": {"tamaño": 20, "tipo_fuente": DEFAULT_FONT_NAME}
            }
        }

            self.create_slide_for_tanda(prs, tanda_number, titulo, titulo_orquesta, subtitulo, genero, canciones, positions_initial)

        # Save the presentation
        try:
            prs.save(output_file)
            os.startfile(output_file)
        except Exception as e:
            messagebox.showerror("Error", f"Failed to save the presentation:\n{e}")
